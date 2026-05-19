"""
최종 PPT 생성:
- 기존 FunETF results + KOFIA fee data 병합
- 합성총보수, 증권거래비용, 상장일/설정일, 자산운용사, 투자위험도
- KOFIA 전체 클래스 보수 테이블 스크린샷 포함
- 두 가지 PPT:
  1. 데이터 요약 PPT (테이블 형식)
  2. 개별 스크린샷 PPT (각 항목별 보수 테이블 이미지)
"""
import json
import re
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

OUTPUT_DIR = Path(r"C:\Users\jswon\Desktop\업무\funetf_scraper\output")
IMG_DIR = OUTPUT_DIR / "images"

# ── 파일 로드 ──
def load_latest_json(pattern):
    files = sorted(OUTPUT_DIR.glob(pattern), reverse=True)
    if not files:
        return []
    with open(files[0], "r", encoding="utf-8") as f:
        return json.load(f), files[0].name

original_data, orig_name = load_latest_json("results_*.json")
kofia_data, kofia_name = load_latest_json("kofia_results_*.json")

print(f"기존 results: {orig_name} ({len(original_data)}개)")
print(f"KOFIA results: {kofia_name} ({len(kofia_data)}개)")

# ── 데이터 병합 ──
orig_map = {item["name"]: item for item in original_data}
kofia_map = {item["name"]: item for item in kofia_data}

merged = []
for item in kofia_data:
    name = item["name"]
    orig = orig_map.get(name, {})

    entry = {
        "name": name,
        "type": item.get("type", orig.get("type", "fund")),
        "listing_date": orig.get("listing_date") or item.get("set_date_kofia"),
        "amc": orig.get("amc") or item.get("amc"),
        "risk_grade": orig.get("risk_grade"),
        "total_fee": item.get("total_fee") or orig.get("total_fee"),
        "trading_cost": item.get("trading_cost"),
        "synthetic_fee": item.get("synthetic_fee"),
        "fee_image_path": item.get("fee_image_path") or orig.get("fee_image_path"),
        "kofia_fund_name": item.get("kofia_rows", [None])[0] if item.get("kofia_rows") else None,
        "std_code": item.get("std_code"),
    }
    merged.append(entry)

print(f"\n병합 결과: {len(merged)}개")
has_fee = sum(1 for m in merged if m.get("synthetic_fee"))
has_tc = sum(1 for m in merged if m.get("trading_cost"))
has_img = sum(1 for m in merged if m.get("fee_image_path"))
print(f"  합성총보수: {has_fee}/{len(merged)}")
print(f"  증권거래비용: {has_tc}/{len(merged)}")
print(f"  이미지: {has_img}/{len(merged)}")

# 병합 결과 저장
ts = datetime.now().strftime("%Y%m%d_%H%M")
merged_file = OUTPUT_DIR / f"merged_results_{ts}.json"
with open(merged_file, "w", encoding="utf-8") as f:
    json.dump(merged, f, ensure_ascii=False, indent=2)
print(f"\n✅ 병합 저장: {merged_file.name}")


def fmt_fee(val):
    """보수 값을 % 형식으로 표시"""
    if not val or val == "-":
        return "-"
    try:
        f = float(val)
        if f == 0:
            return "0%"
        return f"{f:.2f}%"
    except:
        return str(val)


def fmt_date(d):
    if not d:
        return "-"
    s = str(d).replace(".", "").replace("-", "")
    if len(s) == 8:
        return f"{s[:4]}.{s[4:6]}.{s[6:8]}"
    return str(d)


def risk_label(grade):
    m = {"1": "1등급(매우높음)", "2": "2등급(높음)", "3": "3등급(중간)",
         "4": "4등급(낮음)", "5": "5등급(매우낮음)"}
    return m.get(str(grade), str(grade) if grade else "-")


# ═══════════════════════════════════════════════
# PPT 1: 데이터 요약 (테이블 형식)
# ═══════════════════════════════════════════════
def create_summary_ppt():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # 제목 슬라이드
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    def add_text(slide, text, l, t, w, h, bold=False, size=18, color=None, align=PP_ALIGN.LEFT, bg=None):
        from pptx.util import Emu
        txBox = slide.shapes.add_textbox(Emu(int(l*914400)), Emu(int(t*914400)),
                                          Emu(int(w*914400)), Emu(int(h*914400)))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.bold = bold
        run.font.size = Pt(size)
        if color:
            run.font.color.rgb = RGBColor(*color)
        if bg:
            from pptx.oxml.ns import qn
            from lxml import etree
            txBox.fill.solid()
            txBox.fill.fore_color.rgb = RGBColor(*bg)
        return txBox

    # 제목
    add_text(slide, "연금 포트폴리오 ETF/펀드 보수 비용 현황", 0.3, 0.2, 15, 0.6,
             bold=True, size=28, color=(30, 39, 97))

    add_text(slide, f"기준일: 2026년 4월 말 | 출처: 금융투자협회 전자공시(KOFIA)",
             0.3, 0.85, 12, 0.35, size=13, color=(100, 100, 100))

    # ETF 섹션 타이틀
    add_text(slide, "▶ Focus 연금 ETF (11개)", 0.3, 1.35, 7.5, 0.4,
             bold=True, size=16, color=(255, 255, 255), align=PP_ALIGN.LEFT, bg=(30, 39, 97))

    # ETF 헤더
    headers = ["ETF명", "상장일", "자산운용사", "위험등급", "합성총보수", "증권거래비용"]
    col_w = [4.5, 1.3, 2.2, 1.5, 1.5, 1.5]
    col_x = [0.3]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    header_y = 1.8
    for i, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
        add_text(slide, h, x, header_y, w, 0.35, bold=True, size=11,
                 color=(255, 255, 255), align=PP_ALIGN.CENTER, bg=(70, 100, 150))

    # ETF 데이터
    etf_items = [m for m in merged if m["type"] == "etf"]
    row_y = header_y + 0.38
    for i, item in enumerate(etf_items):
        bg = (245, 248, 255) if i % 2 == 0 else (255, 255, 255)
        vals = [
            item["name"][:25],
            fmt_date(item["listing_date"]),
            (item["amc"] or "-")[:10],
            risk_label(item["risk_grade"]),
            fmt_fee(item["synthetic_fee"]),
            fmt_fee(item["trading_cost"]),
        ]
        for j, (val, x, w) in enumerate(zip(vals, col_x, col_w)):
            align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            add_text(slide, val, x, row_y, w, 0.34, size=10, bg=bg, align=align)
        row_y += 0.35

    # 펀드 섹션 (2단 구성)
    fund_items = [m for m in merged if m["type"] == "fund"]
    retire_funds = fund_items[:10]   # 퇴직연금
    pension_funds = fund_items[10:]  # 연금저축

    section_y = row_y + 0.15

    # 퇴직연금 펀드
    add_text(slide, "▶ Focus 퇴직연금 펀드 (10개)", 0.3, section_y, 7.5, 0.38,
             bold=True, size=14, color=(255, 255, 255), bg=(20, 70, 130))
    section_y += 0.42

    for i, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
        add_text(slide, h, x, section_y, w, 0.32, bold=True, size=10,
                 color=(255, 255, 255), align=PP_ALIGN.CENTER, bg=(70, 100, 150))
    section_y += 0.35

    for i, item in enumerate(retire_funds):
        bg = (245, 248, 255) if i % 2 == 0 else (255, 255, 255)
        vals = [
            item["name"][:33],
            fmt_date(item["listing_date"]),
            (item["amc"] or "-")[:8],
            risk_label(item["risk_grade"]),
            fmt_fee(item["synthetic_fee"]),
            fmt_fee(item["trading_cost"]),
        ]
        for j, (val, x, w) in enumerate(zip(vals, col_x, col_w)):
            align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            add_text(slide, val, x, section_y, w, 0.30, size=9, bg=bg, align=align)
        section_y += 0.31

    # 연금저축 펀드 (오른쪽 컬럼)
    col_x2 = [8.5]
    for w in col_w[:-1]:
        col_x2.append(col_x2[-1] + w)

    add_text(slide, "▶ Focus 연금저축 펀드 (10개)", 8.2, 1.35, 7.5, 0.38,
             bold=True, size=14, color=(255, 255, 255), bg=(100, 40, 80))
    sec2_y = 1.8

    for i, (h, x, w) in enumerate(zip(headers, col_x2, col_w)):
        add_text(slide, h, x, sec2_y, w, 0.32, bold=True, size=10,
                 color=(255, 255, 255), align=PP_ALIGN.CENTER, bg=(140, 70, 110))
    sec2_y += 0.35

    for i, item in enumerate(pension_funds):
        bg = (255, 248, 252) if i % 2 == 0 else (255, 255, 255)
        vals = [
            item["name"][:33],
            fmt_date(item["listing_date"]),
            (item["amc"] or "-")[:8],
            risk_label(item["risk_grade"]),
            fmt_fee(item["synthetic_fee"]),
            fmt_fee(item["trading_cost"]),
        ]
        for j, (val, x, w) in enumerate(zip(vals, col_x2, col_w)):
            align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            add_text(slide, val, x, sec2_y, w, 0.30, size=9, bg=bg, align=align)
        sec2_y += 0.31

    out_path = OUTPUT_DIR / f"summary_{ts}.pptx"
    prs.save(str(out_path))
    print(f"\n✅ 요약 PPT 저장: {out_path.name}")
    return out_path


# ═══════════════════════════════════════════════
# PPT 2: 개별 보수 테이블 스크린샷
# ═══════════════════════════════════════════════
def create_screenshot_ppt():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    for idx, item in enumerate(merged):
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)

        # 배경색 설정
        bg = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(0), Inches(16), Inches(9)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(248, 249, 252)
        bg.line.fill.background()

        # 상단 헤더 바
        header_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(16), Inches(1.1))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = RGBColor(30, 39, 97)
        header_bar.line.fill.background()

        # 펀드명
        txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.08), Inches(11), Inches(0.55))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = item["name"]
        run.font.bold = True
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(255, 255, 255)

        # 인덱스 번호
        idx_box = slide.shapes.add_textbox(Inches(14.5), Inches(0.15), Inches(1.2), Inches(0.5))
        tf2 = idx_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        r2 = p2.add_run()
        r2.text = f"{idx+1}/{len(merged)}"
        r2.font.size = Pt(12)
        r2.font.color.rgb = RGBColor(200, 210, 240)

        # 기본 정보 박스 (왼쪽 사이드바)
        info_items = [
            ("상장일/설정일", fmt_date(item.get("listing_date"))),
            ("자산운용사", item.get("amc") or "-"),
            ("투자위험도", risk_label(item.get("risk_grade"))),
            ("합성총보수", fmt_fee(item.get("synthetic_fee"))),
            ("증권거래비용", fmt_fee(item.get("trading_cost"))),
            ("총보수", fmt_fee(item.get("total_fee"))),
        ]

        info_y = 1.25
        for label, val in info_items:
            # 라벨
            lb = slide.shapes.add_textbox(Inches(0.2), Inches(info_y), Inches(1.5), Inches(0.4))
            tf_l = lb.text_frame
            p_l = tf_l.paragraphs[0]
            r_l = p_l.add_run()
            r_l.text = label
            r_l.font.size = Pt(10)
            r_l.font.bold = True
            r_l.font.color.rgb = RGBColor(80, 80, 80)

            # 값
            vb = slide.shapes.add_textbox(Inches(1.75), Inches(info_y), Inches(2.0), Inches(0.4))
            tf_v = vb.text_frame
            p_v = tf_v.paragraphs[0]
            r_v = p_v.add_run()
            r_v.text = str(val)
            r_v.font.size = Pt(12)
            r_v.font.bold = True
            r_v.font.color.rgb = RGBColor(20, 60, 120)

            info_y += 0.55

        # 출처 표시
        src_box = slide.shapes.add_textbox(Inches(0.2), Inches(7.8), Inches(3.5), Inches(0.3))
        tf_s = src_box.text_frame
        p_s = tf_s.paragraphs[0]
        r_s = p_s.add_run()
        r_s.text = "출처: 금융투자협회 전자공시(dis.kofia.or.kr)"
        r_s.font.size = Pt(8)
        r_s.font.color.rgb = RGBColor(150, 150, 150)

        # 스크린샷 이미지 삽입
        img_path = item.get("fee_image_path")
        if img_path and Path(img_path).exists():
            try:
                # 이미지 크기 확인
                with Image.open(img_path) as img:
                    iw, ih = img.size
                    aspect = ih / iw

                # 가용 영역: x=3.8~15.8, y=1.15~8.85 → 약 12인치 × 7.7인치
                max_w = 11.8
                max_h = 7.5
                if aspect * max_w > max_h:
                    h = max_h
                    w = h / aspect
                else:
                    w = max_w
                    h = w * aspect

                # 가운데 정렬
                x = 3.9 + (max_w - w) / 2
                y = 1.2 + (max_h - h) / 2

                slide.shapes.add_picture(img_path,
                                          Inches(x), Inches(y),
                                          Inches(w), Inches(h))
                print(f"  [{idx+1}] ✅ 이미지 삽입: {Path(img_path).name}")
            except Exception as e:
                print(f"  [{idx+1}] ⚠️ 이미지 오류: {e}")
                # 이미지 없음 표시
                no_img = slide.shapes.add_textbox(Inches(6), Inches(4), Inches(6), Inches(1))
                tf = no_img.text_frame
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                r = p.add_run()
                r.text = "⚠️ 이미지 없음"
                r.font.size = Pt(18)
                r.font.color.rgb = RGBColor(200, 100, 100)
        else:
            # 이미지 없음
            no_img = slide.shapes.add_textbox(Inches(6), Inches(4), Inches(6), Inches(1))
            tf = no_img.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = "보수 테이블 이미지 없음"
            r.font.size = Pt(18)
            r.font.color.rgb = RGBColor(180, 180, 180)
            print(f"  [{idx+1}] ❌ 이미지 없음: {item['name'][:30]}")

    out_path = OUTPUT_DIR / f"fee_screenshots_{ts}.pptx"
    prs.save(str(out_path))
    print(f"\n✅ 스크린샷 PPT 저장: {out_path.name}")
    return out_path


# ─── 실행 ───
print("\n" + "=" * 60)
print("PPT 생성 시작")
print("=" * 60)

ppt1 = create_summary_ppt()
ppt2 = create_screenshot_ppt()

print("\n" + "=" * 60)
print("✅ 완료!")
print("=" * 60)
print(f"  1. 요약 PPT: {ppt1.name}")
print(f"  2. 스크린샷 PPT: {ppt2.name}")
print(f"\n저장 위치: {OUTPUT_DIR}")
