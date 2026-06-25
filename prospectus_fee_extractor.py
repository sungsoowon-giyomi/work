"""
funetf 투자설명서 다운로드 + 보수 표 캡처 + PPT 생성

방식: 버튼 직접 클릭 → 다운로드 인터셉트 (URL 패턴 불필요)
"""
import asyncio
import json
import re
import requests
from io import BytesIO
from pathlib import Path
from datetime import datetime

import fitz          # PyMuPDF
import pdfplumber
from PIL import Image
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# ── 설정 ──────────────────────────────────────────────────────────────────
EMAIL    = "sungsoowon45@gmail.com"
PASSWORD = "!sungsoo0405"

OUT_DIR  = Path("output/prospectus")
PDF_DIR  = OUT_DIR / "pdfs"
IMG_DIR  = OUT_DIR / "imgs"
OUT_DIR.mkdir(parents=True, exist_ok=True)
PDF_DIR.mkdir(exist_ok=True)
IMG_DIR.mkdir(exist_ok=True)

TIMESTAMP = datetime.now().strftime("%Y%m%d_%H%M")

# ── ETF 11개 ──────────────────────────────────────────────────────────────
TEST_ITEMS = [
    {"name": "KODEX 미국S&P500",              "type": "etf", "sotCd": "379800"},
    {"name": "TIME 글로벌AI인공지능액티브",      "type": "etf", "sotCd": "456600"},
    {"name": "KODEX AI전력핵심설비",            "type": "etf", "sotCd": "487240"},
    {"name": "ACE 글로벌반도체TOP4 Plus",       "type": "etf", "sotCd": "446770"},
    {"name": "RISE 대형고배당10TR",             "type": "etf", "sotCd": "315960"},
    {"name": "KODEX 머니마켓액티브",             "type": "etf", "sotCd": "488770"},
    {"name": "RISE 미국단기투자등급회사채액티브",  "type": "etf", "sotCd": "437350"},
    {"name": "ACE KRX금현물",                  "type": "etf", "sotCd": "411060"},
    {"name": "KODEX 미국부동산리츠(H)",          "type": "etf", "sotCd": "352560"},
    {"name": "TIGER 글로벌멀티에셋TIF액티브",     "type": "etf", "sotCd": "440340"},
    {"name": "KODEX TRF3070",                 "type": "etf", "sotCd": "329650"},
]


# ── 1. 로그인 ─────────────────────────────────────────────────────────────
async def login(page):
    await page.goto("https://www.funetf.co.kr/auth")
    await page.wait_for_load_state("domcontentloaded")
    await page.fill('input[name="username"]', EMAIL)
    await page.fill('input[name="password"]', PASSWORD)
    await page.locator('button:has-text("로그인")').last.click()
    await page.wait_for_url("https://www.funetf.co.kr/", timeout=15000)
    print("  ✅ 로그인 성공")


# ── 2. 투자설명서 버튼 클릭 → 다운로드 인터셉트 ──────────────────────────
async def download_prospectus_by_click(page, context, sotCd: str, save_path: Path) -> bool:
    """
    funetf 투자설명서 다운로드:
    검색 → 첫 번째 상품 클릭 → '투자설명서' 레이블 아래 다운로드 버튼 클릭
    """

    # ── Step 1: 검색 페이지로 이동 ──────────────────────────────────────
    search_url = f"https://www.funetf.co.kr/search?schVal={sotCd}"
    print(f"  [1] 검색: {search_url}")
    await page.goto(search_url)
    await page.wait_for_load_state("networkidle", timeout=25000)
    await asyncio.sleep(2)

    # ── Step 2: 첫 번째 검색 결과 클릭 ──────────────────────────────────
    print("  [2] 첫 번째 상품 클릭...")

    # 검색 결과 첫 번째 링크 클릭 (여러 셀렉터 시도)
    clicked_url = None
    first_link = await page.evaluate(f'''() => {{
        // sotCd 포함된 링크 우선
        const byTicker = document.querySelector('a[href*="{sotCd}"]');
        if (byTicker) return byTicker.href;
        // 검색 결과 목록의 첫 번째 링크
        const listLinks = document.querySelectorAll(
            '.search-result a, .result-list a, .list-wrap a, ul li a, .item a'
        );
        for (const a of listLinks) {{
            if (a.href && a.href.includes('view')) return a.href;
        }}
        // fallback: view 포함된 첫 번째 링크
        for (const a of document.querySelectorAll('a')) {{
            if (a.href.includes('/product/') || a.href.includes('/view/')) return a.href;
        }}
        return null;
    }}''')

    if first_link:
        print(f"  → {first_link}")
        await page.goto(first_link)
        await page.wait_for_load_state("networkidle", timeout=20000)
        await asyncio.sleep(3)
        clicked_url = page.url
    else:
        print("  ⚠️ 링크 못 찾음, 페이지 내 링크 목록:")
        all_links = await page.evaluate('''() =>
            Array.from(document.querySelectorAll("a"))
                .map(a => ({text: a.textContent?.trim().substring(0,40), href: a.href}))
                .filter(l => l.href && l.href.includes("funetf"))
                .slice(0, 10)
        ''')
        for l in all_links:
            print(f"    {l['text']} → {l['href']}")
        return False

    print(f"  ✅ 상품 페이지: {clicked_url}")

    # ── Step 3: '투자설명서' 레이블 아래 다운로드 버튼 클릭 ────────────
    print("  [3] '투자설명서' 다운로드 버튼 탐색...")

    # 페이지 DOM 구조 파악: '투자설명서' 텍스트 주변 요소 확인
    doc_area = await page.evaluate('''() => {
        const result = [];
        // '투자설명서' 텍스트를 포함한 모든 요소 찾기
        document.querySelectorAll("*").forEach(el => {
            const text = el.childNodes.length === 1 &&
                         el.childNodes[0].nodeType === 3 &&
                         el.textContent?.trim();
            if (text === "투자설명서" || text === "투자설명서 ") {
                // 부모와 형제 요소 확인
                const parent = el.parentElement;
                const siblings = parent ? Array.from(parent.children).map(c => ({
                    tag: c.tagName,
                    text: c.textContent?.trim().substring(0, 60),
                    href: c.getAttribute("href") || "",
                    onclick: (c.getAttribute("onclick") || "").substring(0, 100),
                    class: c.className?.substring(0, 60) || ""
                })) : [];
                result.push({
                    tag: el.tagName,
                    class: el.className || "",
                    parentTag: parent?.tagName || "",
                    parentClass: parent?.className?.substring(0,60) || "",
                    siblings
                });
            }
        });
        return result;
    }''')

    print(f"  '투자설명서' 요소 {len(doc_area)}개 발견:")
    for d in doc_area[:5]:
        print(f"    [{d['tag']}] class={d['class']} / parent=[{d['parentTag']}] class={d['parentClass']}")
        for s in d['siblings']:
            print(f"      sibling [{s['tag']}] '{s['text'][:40]}' href={s['href'][:60]}")

    # '투자설명서' 근처 다운로드 버튼 클릭 (JavaScript로 찾아서 클릭)
    download_btn_info = await page.evaluate('''() => {
        // '투자설명서' 텍스트를 포함한 요소 순회
        const allEls = Array.from(document.querySelectorAll("*"));
        for (const el of allEls) {
            if (el.textContent?.trim() !== "투자설명서") continue;

            // 같은 부모(컨테이너) 안에서 다운로드 버튼 찾기
            const container = el.parentElement;
            if (!container) continue;

            // 컨테이너 내부 + 이후 형제에서 클릭 가능한 요소 찾기
            const searchScope = [container, container.parentElement, container.parentElement?.parentElement];
            for (const scope of searchScope) {
                if (!scope) continue;
                const btns = scope.querySelectorAll("a, button");
                for (const btn of btns) {
                    const t = btn.textContent?.trim();
                    const h = btn.getAttribute("href") || "";
                    const oc = btn.getAttribute("onclick") || "";
                    if (t === "투자설명서") continue; // 레이블 자신 제외
                    if (h.includes(".pdf") || h.includes("invest") || h.includes("download") ||
                        oc.includes("download") || oc.includes("file") ||
                        t.includes("다운") || t.includes("download") || t.includes("DOWN")) {
                        return {
                            found: true,
                            tag: btn.tagName,
                            text: t,
                            href: h,
                            onclick: oc.substring(0, 200)
                        };
                    }
                }
            }
        }
        return {found: false};
    }''')

    print(f"  다운로드 버튼: {download_btn_info}")

    # DOM 구조 확인: div.prd-11-dcmt-item > p.prd-11-dcmt-title("투자설명서") + BUTTON("다운로드")
    # '투자설명서' 카드 안의 '다운로드' 버튼을 정확히 타겟팅
    dl_selector = 'div.prd-11-dcmt-item:has(p.prd-11-dcmt-title:text("투자설명서")) button:has-text("다운로드")'

    try:
        btn = page.locator(dl_selector)
        count = await btn.count()
        print(f"  다운로드 버튼 발견: {count}개 ({dl_selector})")

        if count > 0:
            async with page.expect_download(timeout=30000) as dl_info:
                await btn.first.click()
            dl = await dl_info.value
            await dl.save_as(str(save_path))
            print(f"  ✅ 다운로드 성공: {save_path.name} ({save_path.stat().st_size:,} bytes)")
            return True
    except Exception as e:
        print(f"  다운로드 버튼 클릭 오류: {e}")

    # fallback: 클래스 없이 텍스트만으로 찾기
    print("  fallback 시도...")
    try:
        async with page.expect_download(timeout=20000) as dl_info:
            # '투자설명서' 텍스트를 가진 p 태그의 부모 div에서 '다운로드' 버튼 클릭
            await page.evaluate('''() => {
                const titles = Array.from(document.querySelectorAll("p, span, div"));
                for (const el of titles) {
                    if (el.textContent?.trim() === "투자설명서" && el.children.length === 0) {
                        // 부모 컨테이너에서 '다운로드' 버튼 찾기
                        const container = el.closest("div");
                        if (container) {
                            const btns = container.querySelectorAll("button");
                            for (const btn of btns) {
                                if (btn.textContent?.trim().includes("다운로드")) {
                                    btn.click();
                                    return;
                                }
                            }
                        }
                    }
                }
            }''')
        dl = await dl_info.value
        await dl.save_as(str(save_path))
        print(f"  ✅ fallback 다운로드 성공: {save_path.name}")
        return True
    except Exception as e:
        print(f"  fallback 실패: {e}")

    return False


# ── 3. '(1)투자신탁 관련보수 등' 페이지 찾기 ─────────────────────────────
def find_fee_page(pdf_path: Path) -> int | None:
    # 섹션 헤더 키워드
    keywords_section = [
        "집합투자기구에 부과되는 보수",
        "(1)투자신탁", "(1) 투자신탁",
        "투자신탁 관련보수",
        "투자신탁 관련 보수",  # KODEX TRF3070: "(1) 투자신탁 관련 보수등"
    ]
    # 보수 수치가 있는 페이지 키워드 (스페이스 있/없 모두 허용)
    keywords_fee = [
        "증권거래비용", "증권 거래비용",
        "총보수", "총 보수",
        "합성총보수", "합성 총보수",
    ]

    fallback_pg = None
    try:
        with pdfplumber.open(pdf_path) as pdf:
            total = len(pdf.pages)
            print(f"  PDF 총 {total}페이지")
            for pg_num, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                has_section = any(kw in text for kw in keywords_section)
                has_fee     = any(kw in text for kw in keywords_fee)

                # 1차: 같은 페이지에 섹션 헤더 + 보수 키워드
                if has_section and has_fee:
                    print(f"  → 보수 섹션(1차): {pg_num+1}페이지")
                    return pg_num

                # 2차: 섹션 헤더는 이 페이지, 보수 수치는 다음 페이지
                if has_section and pg_num + 1 < total:
                    next_text = pdf.pages[pg_num + 1].extract_text() or ""
                    if any(kw in next_text for kw in keywords_fee):
                        print(f"  → 보수 섹션(섹션+다음페이지): {pg_num+1}페이지")
                        return pg_num

                # 3차 fallback 후보: 보수 관련 키워드 다수 포함
                if fallback_pg is None:
                    has_synth = any(k in text for k in ["합성총보수", "합성 총보수", "총보수·비용", "총 보수·비용"])
                    has_trade = any(k in text for k in ["증권거래비용", "증권 거래비용"])
                    if has_synth and has_trade:
                        fallback_pg = pg_num

        if fallback_pg is not None:
            print(f"  → 보수 섹션(fallback): {fallback_pg+1}페이지")
            return fallback_pg
    except Exception as e:
        print(f"  페이지 탐색 오류: {e}")
    return None


# ── 4. 표 이미지 캡처 + 크롭 ─────────────────────────────────────────────
def capture_fee_table(pdf_path: Path, pg_num: int, out_img: Path, zoom: float = 2.5) -> Path:
    """
    캡처 범위:
    - 시작: '집합투자기구에 부과되는 보수 및 비용' 텍스트 위쪽
    - 끝:   표의 마지막 행(지급시기) 아래쪽
    - 다음 페이지까지 표가 이어지면 합쳐서 크롭
    """
    doc = fitz.open(str(pdf_path))

    # ── Step 1: 시작 y좌표 찾기 (PyMuPDF로 텍스트 검색) ──────────────────
    # '집합투자기구에 부과되는 보수' 텍스트의 y 위치 (띄어쓰기 있/없 모두 시도)
    section_keywords = [
        "집합투자기구에 부과되는 보수",   # 일반형
        "집합투자기구에 부과",           # 짧은형
        "집합투자기구에부과되는보수",     # 공백없는형 (ACE/일부 PDF)
        "집합투자기구에부과",             # 공백없는 짧은형
        "나. 집합투자기구",              # prefix 형
    ]
    start_y_pdf = None   # PDF 좌표계 (pt)
    page0 = doc[pg_num]

    for kw in section_keywords:
        hits = page0.search_for(kw)
        if hits:
            start_y_pdf = hits[0].y0
            print(f"  섹션 시작 y: {start_y_pdf:.1f}pt ('{kw}')")
            break

    # ── Step 2: 표 끝 y좌표 찾기 (pdfplumber) ────────────────────────────
    # 이 섹션 표가 다음 페이지까지 이어지는지 확인
    # '지급시기' 행이 다음 페이지에 있거나, '증권거래비용'이 다음 페이지 첫 번째 표에 있을 때만 확장
    pages_needed = [pg_num]
    if pg_num + 1 < len(doc):
        next_text = doc[pg_num + 1].get_text()
        # 다음 페이지 첫 200줄 (본문 상단) 만 확인 — 각주까지 포함하지 않기 위해
        next_lines = next_text.split("\n")[:30]
        next_top = "\n".join(next_lines)
        if "지급시기" in next_top or "증권거래비용" in next_top or "합성총보수" in next_top:
            pages_needed.append(pg_num + 1)

    print(f"  캡처 페이지: {[p+1 for p in pages_needed]}")

    # ── Step 3: 각 페이지 렌더링 ─────────────────────────────────────────
    page_images = []
    page_heights_pdf = []
    for pn in pages_needed:
        p = doc[pn]
        pix = p.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
        page_images.append(Image.open(BytesIO(pix.tobytes("png"))))
        page_heights_pdf.append(p.rect.height)
    doc.close()

    # 이미지 세로 합치기
    if len(page_images) > 1:
        total_h = sum(img.height for img in page_images)
        full_img = Image.new("RGB", (page_images[0].width, total_h), "white")
        y_off = 0
        for img in page_images:
            full_img.paste(img, (0, y_off))
            y_off += img.height
    else:
        full_img = page_images[0]

    img_w = full_img.width
    img_h = full_img.height
    single_page_img_h = page_images[0].height
    single_page_img_w = page_images[0].width

    # ── Step 4: 크롭 좌표 계산 ────────────────────────────────────────────
    try:
        with pdfplumber.open(pdf_path) as pdf:
            pl_page0 = pdf.pages[pg_num]
            pw0 = pl_page0.width    # PDF 페이지 너비 (pt)
            ph0 = pl_page0.height   # PDF 페이지 높이 (pt)

            sx = single_page_img_w / pw0   # x 스케일: PDF pt → 이미지 px
            sy = single_page_img_h / ph0   # y 스케일

            x0_pdf, x1_pdf = None, None
            end_y_img = None

            for i, pn in enumerate(pages_needed):
                pl_page = pdf.pages[pn]
                y_offset_img = sum(page_images[j].height for j in range(i))

                fee_kws_cap  = ["총보수", "총 보수", "증권거래비용", "증권 거래비용",
                                 "거래비용", "합성총보수", "피투자"]
                num_re_cap   = re.compile(r"\d{1,2}\.\d{2,4}")

                for tbl in pl_page.find_tables():
                    ext = tbl.extract()
                    if not ext: continue
                    flat = " ".join(str(c) for row in ext for c in row if c)

                    # 보수 관련 키워드 있는 표만
                    if not any(kw in flat for kw in fee_kws_cap):
                        continue

                    # ★ 핵심: 실제 숫자 데이터가 있는 표만 (각주/예시 표 제외)
                    #   유효한 보수 숫자(0.0001~9.9999)가 표 전체에 2개 이상이면 통과
                    #   (3열 표는 행마다 1개씩, 넓은 표는 한 행에 여러 개)
                    total_valid = 0
                    for row in ext:
                        if not row: continue
                        for cell in row:
                            for m in num_re_cap.findall(str(cell or "")):
                                try:
                                    if 0.0001 <= float(m) <= 9.9999:
                                        total_valid += 1
                                except:
                                    pass
                    if total_valid < 2:
                        continue  # 숫자가 거의 없는 표 (순수 텍스트 각주 등) 제외

                    bbox = tbl.bbox  # (x0, top, x1, bottom) in PDF pt
                    if x0_pdf is None:
                        x0_pdf = bbox[0]
                        x1_pdf = bbox[2]
                    # 표 끝 y (항상 마지막 값으로 갱신)
                    end_y_img = y_offset_img + int(bbox[3] * sy) + 20

            if x0_pdf is not None:
                cx0 = max(0, int(x0_pdf * sx) - 10)
                cx1 = min(img_w, int(x1_pdf * sx) + 10)
                # 시작 y: '집합투자기구에 부과되는 보수' 텍스트 위 10px
                cy0 = max(0, int(start_y_pdf * sy) - 10) if start_y_pdf else 0
                cy1 = min(img_h, end_y_img or img_h)

                full_img = full_img.crop((cx0, cy0, cx1, cy1))
                print(f"  ✅ 크롭: {cx1-cx0}×{cy1-cy0}px")
            else:
                print("  ⚠️  표 위치 못 찾음 → 전체 페이지 저장")

    except Exception as e:
        print(f"  크롭 오류: {e}")

    full_img.save(str(out_img))
    print(f"  ✅ 이미지 저장: {out_img.name}")
    return out_img


# ── 5. 보수 값 추출 ───────────────────────────────────────────────────────
def extract_fees(pdf_path: Path, pg_num: int) -> dict:
    """
    PDF 표에서 아래 두 값을 추출:
    - 총보수·비용(피투자집합투자기구 보수포함)  → synthetic_fee
    - 증권거래비용                              → trading_cost

    지원하는 표 형식:
    1) 넓은 다열 표 (KODEX/TIME/TIGER 스타일): 헤더에 "피투자" 열, "거래비용" 열
    2) 3열 표 (ACE/KRX 스타일): [구분, 비율, 시기] 행 단위
    3) 순수 텍스트 (RISE/KRX 국내형): "총 보수·비용 X.XXX" 라인
    """
    result = {"synthetic_fee": None, "trading_cost": None}
    num_find = re.compile(r"\d{1,2}\.\d{2,4}")

    # 유효한 보수 % 범위: 0.001% ~ 9.999% (연도 '25.01' 같은 오탐 방지)
    def is_valid_fee(s: str) -> bool:
        try:
            v = float(s)
            return 0.0001 <= v <= 9.9999
        except:
            return False

    pages_to_check = [pg_num, pg_num + 1]

    # 보수 관련 테이블 감지 키워드 (스페이스 포함 변형도 허용)
    FEE_KWS = ["총보수", "총 보수", "증권거래비용", "증권 거래비용", "거래비용", "합성총보수"]

    # ── 방법 1: 표 파싱 ───────────────────────────────────────────────────
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for pn in pages_to_check:
                if pn >= len(pdf.pages): break
                page = pdf.pages[pn]

                for tbl in page.extract_tables():
                    if not tbl: continue
                    flat = " ".join(str(c) for row in tbl for c in row if c)
                    if not any(kw in flat for kw in FEE_KWS):
                        continue

                    print(f"    표 발견 (p.{pn+1}), 행 수: {len(tbl)}")

                    # ── 3열 표 (ACE/KRX 스타일: 구분 | 비율 | 시기) ──
                    max_cols = max((len(r) for r in tbl if r), default=0)
                    if max_cols <= 5:
                        for row in tbl:
                            row_s = [str(c or "").replace("\n", " ").strip() for c in row]
                            if len(row_s) < 2: continue
                            key = row_s[0]
                            val_str = row_s[1] if len(row_s) > 1 else ""
                            nums = [n for n in num_find.findall(val_str) if is_valid_fee(n)]
                            if not nums: continue
                            # 총보수·비용 행 → 합성총보수 역할 (피투자기구 없는 국내 ETF)
                            if any(k in key for k in ["총보수·비용", "총보수ᆞ비용", "총보수･비용",
                                                       "총 보수·비용", "총보수·비용", "총보수비용"]):
                                if result["synthetic_fee"] is None:
                                    result["synthetic_fee"] = nums[0]
                                    print(f"    3열synth: '{key}'={nums[0]}")
                            if any(k in key for k in ["증권거래비용", "증권 거래비용"]):
                                if result["trading_cost"] is None:
                                    result["trading_cost"] = nums[0]
                                    print(f"    3열trade: '{key}'={nums[0]}")
                        if result["synthetic_fee"] and result["trading_cost"]:
                            return result
                        continue

                    # ── 넓은 다열 표 ──────────────────────────────────────
                    # 헤더에서 열 인덱스 탐지
                    col_synth = None   # 피투자집합투자기구 보수 포함 열
                    col_trade = None   # 증권거래비용 열
                    for row in tbl:
                        row_s = [str(c or "").replace("\n", " ").strip() for c in row]
                        for i, cell in enumerate(row_s):
                            if "피투자" in cell or "보수포함" in cell:
                                col_synth = i
                            if "증권거래비용" in cell or "거래비용" in cell:
                                col_trade = i

                    print(f"    열 인덱스: col_synth={col_synth}, col_trade={col_trade}")

                    for row in tbl:
                        row_s = [str(c or "").replace("\n", " ").strip() for c in row]
                        flat_valid = [n for c in row_s for n in num_find.findall(c) if is_valid_fee(n)]

                        if len(flat_valid) < 3:
                            continue
                        print(f"    데이터 행 숫자(유효): {flat_valid}")

                        # 우선순위 1: 두 열 인덱스 모두 알고, 값이 다름
                        if col_synth is not None and col_trade is not None:
                            if col_synth < len(row_s) and col_trade < len(row_s):
                                sn = [n for n in num_find.findall(row_s[col_synth]) if is_valid_fee(n)]
                                tn = [n for n in num_find.findall(row_s[col_trade]) if is_valid_fee(n)]
                                if sn and tn:
                                    result["synthetic_fee"] = sn[0]
                                    result["trading_cost"]  = tn[0]
                                    print(f"    ✅ 두열인덱스: synth={sn[0]}, trade={tn[0]}")
                                    return result

                        # 우선순위 2: col_synth만 알 때 (TIGER TIF: 증권거래비용 열 없음)
                        if col_synth is not None and col_trade is None:
                            if col_synth < len(row_s):
                                sn = [n for n in num_find.findall(row_s[col_synth]) if is_valid_fee(n)]
                                if sn and result["synthetic_fee"] is None:
                                    result["synthetic_fee"] = sn[0]
                                    print(f"    ✅ col_synth: synth={sn[0]}")
                            continue  # 이 행에서 trading_cost 없으면 다음 행으로

                        # 우선순위 3: col_trade만 알 때
                        if col_trade is not None and col_synth is None:
                            if col_trade < len(row_s):
                                tn = [n for n in num_find.findall(row_s[col_trade]) if is_valid_fee(n)]
                                if tn:
                                    result["trading_cost"] = tn[0]
                                    # 합성총보수는 바로 앞 열에서 추정
                                    if len(flat_valid) >= 2:
                                        result["synthetic_fee"] = flat_valid[-2] if flat_valid[-1] == tn[0] else flat_valid[-1]
                                    print(f"    ✅ col_trade만: synth={result['synthetic_fee']}, trade={tn[0]}")
                                    return result

                        # 우선순위 4: 마지막 두 숫자 (값이 다를 때만)
                        if len(flat_valid) >= 2 and flat_valid[-2] != flat_valid[-1]:
                            result["synthetic_fee"] = flat_valid[-2]
                            result["trading_cost"]  = flat_valid[-1]
                            print(f"    ✅ 마지막2: synth={flat_valid[-2]}, trade={flat_valid[-1]}")
                            return result

                    if result["synthetic_fee"] and result["trading_cost"]:
                        return result
    except Exception as e:
        print(f"    표 파싱 오류: {e}")

    # ── 방법 2: 텍스트 라인 파싱 ─────────────────────────────────────────
    print("    텍스트 방식으로 fallback...")
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for pn in pages_to_check:
                if pn >= len(pdf.pages): break
                text = pdf.pages[pn].extract_text() or ""
                lines = text.split("\n")

                for line in lines:
                    l = line.strip()
                    if not l:
                        continue
                    nums_valid = [n for n in num_find.findall(l) if is_valid_fee(n)]
                    if not nums_valid:
                        continue

                    # 합성총보수·피투자포함 행
                    if any(k in l for k in ["합성총보수", "합성 총보수"]):
                        if result["synthetic_fee"] is None:
                            result["synthetic_fee"] = nums_valid[-1]
                            print(f"    합성총보수 행: {l[:80]} → {nums_valid[-1]}")
                    elif "피투자" in l and result["synthetic_fee"] is None:
                        result["synthetic_fee"] = nums_valid[-1]
                        print(f"    피투자 행: {l[:80]} → {nums_valid[-1]}")
                    # 총보수·비용 (국내 ETF의 최종 합계 — 피투자 없음)
                    elif any(k in l for k in ["총 보수·비용", "총보수·비용", "총보수ᆞ비용",
                                               "총보수·비용", "총보수･비용", "총보수비용"]) \
                            and result["synthetic_fee"] is None:
                        result["synthetic_fee"] = nums_valid[-1]
                        print(f"    총보수비용 행: {l[:80]} → {nums_valid[-1]}")

                    # 증권거래비용
                    if any(k in l for k in ["증권거래비용", "증권 거래비용"]):
                        if result["trading_cost"] is None:
                            result["trading_cost"] = nums_valid[-1]
                            print(f"    거래비용 행: {l[:80]} → {nums_valid[-1]}")

                if result["synthetic_fee"] and result["trading_cost"]:
                    return result

                # 숫자 8개+ 연속 행 (한 줄에 표 전체가 들어간 경우)
                for line in lines:
                    tokens = line.strip().split()
                    valid = [t for t in tokens if num_find.fullmatch(t) and is_valid_fee(t)]
                    if len(valid) >= 8:
                        print(f"    숫자8+행: {valid}")
                        result["synthetic_fee"] = valid[-2]
                        result["trading_cost"]  = valid[-1]
                        return result
    except Exception as e:
        print(f"    텍스트 파싱 오류: {e}")

    # ── 방법 3: 워드 스캔 ────────────────────────────────────────────────
    print("    워드 스캔 fallback...")
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for pn in pages_to_check:
                if pn >= len(pdf.pages): break
                words = pdf.pages[pn].extract_words()
                wt = [w["text"] for w in words]
                for i, w in enumerate(wt):
                    if any(k in w for k in ["합성총보수", "피투자"]) and result["synthetic_fee"] is None:
                        for j in range(i+1, min(i+15, len(wt))):
                            ns = [n for n in num_find.findall(wt[j]) if is_valid_fee(n)]
                            if ns:
                                result["synthetic_fee"] = ns[0]
                                print(f"    워드synth: {wt[j]} → {ns[0]}")
                                break
                    if any(k in w for k in ["증권거래비용", "거래비용"]) and result["trading_cost"] is None:
                        for j in range(i+1, min(i+15, len(wt))):
                            ns = [n for n in num_find.findall(wt[j]) if is_valid_fee(n)]
                            if ns:
                                result["trading_cost"] = ns[0]
                                print(f"    워드trade: {wt[j]} → {ns[0]}")
                                break
                if result["synthetic_fee"] and result["trading_cost"]:
                    return result
    except Exception as e:
        print(f"    워드스캔 오류: {e}")

    return result


# ── 6. PPT 생성 ──────────────────────────────────────────────────────────
def set_cell_text(cell, text, font_size=11, bold=False,
                  font_name="맑은 고딕",
                  align=None,
                  bg_color=None,
                  text_color=RGBColor(0, 0, 0)):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor as RC
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.clear()
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.east_asian_font_name = font_name
    run.font.color.rgb = text_color

    if bg_color:
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
        srgbClr  = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", str(bg_color))


def add_summary_slide(prs, results: list):
    """첫 슬라이드: 상품명 / 티커 / 총보수·비용(피투자포함) / 증권거래비용 표"""
    from pptx.enum.text import PP_ALIGN
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # 제목
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.55))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "ETF 보수 요약"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.name = "맑은 고딕"
    run.font.east_asian_font_name = "맑은 고딕"
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # 표: 헤더 1행 + 데이터 N행
    cols = 4
    rows = 1 + len(results)
    tbl = slide.shapes.add_table(rows, cols,
                                  Inches(0.4), Inches(0.85),
                                  Inches(12.5), Inches(0.45 * rows)).table

    # 열 너비
    tbl.columns[0].width = Inches(5.5)   # 상품명
    tbl.columns[1].width = Inches(1.2)   # 티커
    tbl.columns[2].width = Inches(3.0)   # 총보수·비용(피투자포함)
    tbl.columns[3].width = Inches(2.0)   # 증권거래비용

    hdr_bg  = RGBColor(0x1F, 0x49, 0x7D)
    hdr_fg  = RGBColor(0xFF, 0xFF, 0xFF)
    row_bg1 = RGBColor(0xF2, 0xF6, 0xFF)
    row_bg2 = RGBColor(0xFF, 0xFF, 0xFF)

    headers = ["상품명", "티커", "총보수·비용\n(피투자집합투자기구 보수포함)", "증권거래비용"]
    for ci, h in enumerate(headers):
        set_cell_text(tbl.cell(0, ci), h,
                      font_size=10, bold=True,
                      align=PP_ALIGN.CENTER,
                      bg_color=hdr_bg, text_color=hdr_fg)

    for ri, r in enumerate(results, start=1):
        bg = row_bg1 if ri % 2 == 1 else row_bg2
        ticker = next((item["sotCd"] for item in TEST_ITEMS if item["name"] == r["name"]), "")
        sf = f"{r.get('synthetic_fee', 'N/A')}%" if r.get('synthetic_fee') else "N/A"
        tc = f"{r.get('trading_cost',  'N/A')}%" if r.get('trading_cost')  else "N/A"
        row_data = [r["name"], ticker, sf, tc]
        aligns   = [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.CENTER]
        for ci, (val, aln) in enumerate(zip(row_data, aligns)):
            set_cell_text(tbl.cell(ri, ci), val,
                          font_size=10, align=aln, bg_color=bg)


def create_ppt(results: list, out_path: Path):
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── 첫 슬라이드: 요약 표 ──
    add_summary_slide(prs, results)

    for r in results:
        slide = prs.slides.add_slide(blank)

        def add_text(slide, text, left, top, width, height, size, bold=False, color=RGBColor(0,0,0)):
            tb = slide.shapes.add_textbox(left, top, width, height)
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = "맑은 고딕"          # ← 한국어 폰트 명시
            run.font.east_asian_font_name = "맑은 고딕"
            return tb

        # ETF 이름
        add_text(slide, r["name"],
                 Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.5),
                 size=18, bold=True, color=RGBColor(0x1F, 0x49, 0x7D))

        # 보수 수치
        sf = f"{r['synthetic_fee']}%" if r.get('synthetic_fee') else "N/A"
        tc = f"{r['trading_cost']}%"  if r.get('trading_cost')  else "N/A"
        add_text(slide,
                 f"총보수·비용(피투자집합투자기구 보수포함): {sf}    증권거래비용: {tc}",
                 Inches(0.3), Inches(0.65), Inches(12.7), Inches(0.35),
                 size=11, color=RGBColor(0x55, 0x55, 0x55))

        # 표 이미지
        img_path = r.get("img_path")
        if img_path and Path(img_path).exists():
            img = Image.open(img_path)
            iw, ih = img.size
            max_w, max_h = Inches(12.7), Inches(6.3)
            ratio   = min(max_w / iw, max_h / ih)
            disp_w  = iw * ratio
            disp_h  = ih * ratio
            left    = Inches(0.3) + (max_w - disp_w) / 2
            slide.shapes.add_picture(img_path, left, Inches(1.05), disp_w, disp_h)

    prs.save(str(out_path))
    print(f"✅ PPT 저장: {out_path}")


# ── 메인 ─────────────────────────────────────────────────────────────────
async def main():
    print("=" * 60)
    print("투자설명서 보수 표 추출기 (ETF 3개 테스트)")
    print("=" * 60)

    results = []

    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        context = await browser.new_context(
            user_agent='Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 Chrome/120.0.0.0 Safari/537.36',
            accept_downloads=True,
        )
        page = await context.new_page()

        print("\n[로그인]")
        await login(page)

        for item in TEST_ITEMS:
            name  = item["name"]
            sotCd = item["sotCd"]
            safe  = name.replace(" ", "_").replace("/", "_")
            pdf_path = PDF_DIR / f"{safe}.pdf"
            img_path = IMG_DIR / f"{safe}_fee.png"

            print(f"\n{'='*50}")
            print(f"▶ {name} ({sotCd})")
            print(f"{'='*50}")

            r = {"name": name, "img_path": None, "synthetic_fee": None, "trading_cost": None}

            # 이미 다운로드된 PDF 재사용
            if pdf_path.exists() and pdf_path.stat().st_size > 5000:
                print(f"  PDF 이미 존재 → 재사용: {pdf_path.name}")
            else:
                print(f"  [1] PDF 다운로드: funetf.co.kr/search?schVal={sotCd}")
                ok = await download_prospectus_by_click(page, context, sotCd, pdf_path)
                if not ok:
                    print("  ❌ 다운로드 실패")
                    results.append(r)
                    continue

            print("  [2] 보수 섹션 페이지 탐색...")
            pg = find_fee_page(pdf_path)
            if pg is None:
                print("  ❌ 보수 섹션 없음")
                results.append(r)
                continue

            print("  [3] 표 이미지 캡처...")
            saved = capture_fee_table(pdf_path, pg, img_path)
            r["img_path"] = str(saved)

            print("  [4] 보수 값 추출...")
            fees = extract_fees(pdf_path, pg)
            r["synthetic_fee"] = fees["synthetic_fee"]
            r["trading_cost"]  = fees["trading_cost"]
            print(f"  ▶ 합성총보수:   {fees['synthetic_fee']}%")
            print(f"  ▶ 증권거래비용: {fees['trading_cost']}%")

            results.append(r)

        await browser.close()

    # PPT 생성
    ppt_path = OUT_DIR / f"fee_tables_{TIMESTAMP}.pptx"
    print(f"\n[PPT 생성]")
    create_ppt(results, ppt_path)

    # 결과 요약
    print("\n" + "=" * 60)
    print(f"{'ETF명':<35} {'합성총보수':>8} {'증권거래비용':>10}")
    print("-" * 60)
    for r in results:
        sf = f"{r['synthetic_fee']}%" if r['synthetic_fee'] else "N/A"
        tc = f"{r['trading_cost']}%" if r['trading_cost'] else "N/A"
        print(f"{r['name']:<35} {sf:>8} {tc:>10}")

    # 저장
    (OUT_DIR / f"fee_results_{TIMESTAMP}.json").write_text(
        json.dumps(results, ensure_ascii=False, indent=2), encoding="utf-8"
    )


asyncio.run(main())
