"""
우주항공 ETF 보고서 생성기
- 한국 ETF (4개): funetf.co.kr → PDF 다운 → 보수 표 캡처
- 미국 ETF (4개): etfcheck.co.kr → 기본정보 캡처 + 설정일/시가총액/TER 추출
- PPT: 요약 슬라이드 + 개별 슬라이드
"""
import asyncio
import json
import re
import requests
from io import BytesIO
from pathlib import Path
from datetime import datetime

import fitz
import pdfplumber
from PIL import Image
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 설정 ──────────────────────────────────────────────────────────────────
EMAIL    = "sungsoowon45@gmail.com"
PASSWORD = "!sungsoo0405"

OUT_DIR  = Path("output/space_etf")
PDF_DIR  = OUT_DIR / "pdfs"
IMG_DIR  = OUT_DIR / "imgs"
OUT_DIR.mkdir(parents=True, exist_ok=True)
PDF_DIR.mkdir(exist_ok=True)
IMG_DIR.mkdir(exist_ok=True)

TIMESTAMP = datetime.now().strftime("%Y%m%d_%H%M")

# ── 한국 ETF (신규 알파벳 코드, 'A' 제거) ────────────────────────────────
KR_ITEMS = [
    {"name": "TIGER 미국우주테크",      "sotCd": "0183J0"},
    {"name": "KODEX 미국우주항공",      "sotCd": "0167Z0"},
    {"name": "1Q 미국우주항공테크",     "sotCd": "0131V0"},
    {"name": "ACE 미국우주테크액티브",  "sotCd": "0180V0"},
]

# ── 미국 ETF (etfcheck.co.kr ID) ─────────────────────────────────────────
US_ITEMS = [
    {"name": "Tema Space Innovators ETF",  "ticker": "NASA", "etfcheck_id": "F00001TIVN"},
    {"name": "Procure Space ETF",           "ticker": "UFO",  "etfcheck_id": "F00000ZZBI"},
    {"name": "ARK Space Exploration ETF",  "ticker": "ARKX", "etfcheck_id": "F000016CI6"},
    {"name": "N/A (etfcheck 미등재)",       "ticker": "SPCF", "etfcheck_id": None},
]


# ═══════════════════════════════════════════════════════════════════════════
# ── SECTION 1: 한국 ETF (funetf → PDF → 보수 추출) ──────────────────────
# ═══════════════════════════════════════════════════════════════════════════

async def login(page):
    await page.goto("https://www.funetf.co.kr/auth")
    await page.wait_for_load_state("domcontentloaded")
    await page.fill('input[name="username"]', EMAIL)
    await page.fill('input[name="password"]', PASSWORD)
    await page.locator('button:has-text("로그인")').last.click()
    await page.wait_for_url("https://www.funetf.co.kr/", timeout=15000)
    print("  ✅ 로그인 성공")


async def download_kr_pdf(page, context, sotCd: str, save_path: Path) -> bool:
    """funetf.co.kr 검색 → 상품 페이지 → 투자설명서 다운로드"""

    # Step 1: 검색
    search_url = f"https://www.funetf.co.kr/search?schVal={sotCd}"
    print(f"  [1] 검색: {search_url}")
    await page.goto(search_url)
    await page.wait_for_load_state("networkidle", timeout=25000)
    await asyncio.sleep(2)

    # Step 2: 첫 번째 상품 클릭
    first_link = await page.evaluate(f'''() => {{
        const byTicker = document.querySelector('a[href*="{sotCd}"]');
        if (byTicker) return byTicker.href;
        for (const a of document.querySelectorAll('a')) {{
            if (a.href.includes('/product/') || a.href.includes('/view/')) return a.href;
        }}
        return null;
    }}''')

    if not first_link:
        print(f"  ❌ 상품 링크 없음")
        return False

    print(f"  [2] 상품 페이지: {first_link}")
    await page.goto(first_link)
    await page.wait_for_load_state("networkidle", timeout=20000)
    await asyncio.sleep(3)

    # Step 3: 투자설명서 다운로드 버튼 클릭
    dl_selector = 'div.prd-11-dcmt-item:has(p.prd-11-dcmt-title:text("투자설명서")) button:has-text("다운로드")'
    try:
        btn = page.locator(dl_selector)
        if await btn.count() > 0:
            async with page.expect_download(timeout=30000) as dl_info:
                await btn.first.click()
            dl = await dl_info.value
            await dl.save_as(str(save_path))
            print(f"  ✅ PDF 저장: {save_path.name} ({save_path.stat().st_size:,} bytes)")
            return True
    except Exception as e:
        print(f"  버튼 클릭 오류: {e}")

    # fallback: data-fileurl 방식
    try:
        fileurl = await page.evaluate('''() => {
            const items = document.querySelectorAll(".prd-11-dcmt-item");
            for (const item of items) {
                const title = item.querySelector(".prd-11-dcmt-title");
                if (title && title.textContent?.trim() === "투자설명서") {
                    const btn = item.querySelector("button.download");
                    return btn ? btn.getAttribute("data-fileurl") : null;
                }
            }
            return null;
        }''')
        if fileurl:
            full_url = f"https://www.funetf.co.kr{fileurl}"
            print(f"  [fallback] 직접 다운로드: {full_url}")
            # requests로 다운로드 (쿠키 필요 없는 경우)
            resp = requests.get(full_url, timeout=30)
            if resp.ok and len(resp.content) > 1000:
                save_path.write_bytes(resp.content)
                print(f"  ✅ 저장: {save_path.name} ({len(resp.content):,} bytes)")
                return True
    except Exception as e:
        print(f"  fallback 오류: {e}")

    return False


def find_fee_page(pdf_path: Path) -> int | None:
    keywords_section = [
        "집합투자기구에 부과되는 보수",
        "(1)투자신탁", "(1) 투자신탁",
        "투자신탁 관련보수",
        "투자신탁 관련 보수",
    ]
    keywords_fee = [
        "증권거래비용", "증권 거래비용",
        "총보수", "총 보수",
        "합성총보수", "합성 총보수",
    ]
    candidates = []   # (pg_num, ) — 섹션+보수 키워드 모두 있는 페이지
    fallback_pg = None
    try:
        with pdfplumber.open(pdf_path) as pdf:
            total = len(pdf.pages)
            print(f"  PDF 총 {total}페이지")
            for pg_num, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                has_section = any(kw in text for kw in keywords_section)
                has_fee     = any(kw in text for kw in keywords_fee)
                if has_section and has_fee:
                    print(f"  → 보수 섹션 후보: {pg_num+1}페이지")
                    candidates.append(pg_num)
                elif has_section and pg_num + 1 < total:
                    next_text = pdf.pages[pg_num + 1].extract_text() or ""
                    if any(kw in next_text for kw in keywords_fee):
                        print(f"  → 보수 섹션 후보(섹션+다음페이지): {pg_num+1}페이지")
                        candidates.append(pg_num)
                if fallback_pg is None:
                    has_synth = any(k in text for k in ["합성총보수", "합성 총보수", "총보수·비용", "총보수·비",
                                                        "총보수･비용", "총보수･비", "총 보수·비용"])
                    has_trade = any(k in text for k in ["증권거래비용", "증권 거래비용"])
                    if has_synth and has_trade:
                        fallback_pg = pg_num

        if candidates:
            # 간이투자설명서(초반)보다 상세 투자설명서(후반) 섹션을 우선 선택
            pg = candidates[-1]
            print(f"  → 선택된 보수 섹션: {pg+1}페이지 (후보 {len(candidates)}개 중 마지막)")
            return pg
        if fallback_pg is not None:
            print(f"  → 보수 섹션(fallback): {fallback_pg+1}페이지")
            return fallback_pg
    except Exception as e:
        print(f"  페이지 탐색 오류: {e}")
    return None


def capture_fee_table(pdf_path: Path, pg_num: int, out_img: Path, zoom: float = 2.5) -> Path:
    doc = fitz.open(str(pdf_path))
    section_keywords = [
        "집합투자기구에 부과되는 보수",
        "집합투자기구에 부과",
        "집합투자기구에부과되는보수",
        "집합투자기구에부과",
        "나. 집합투자기구",
    ]
    start_y_pdf = None
    page0 = doc[pg_num]

    for kw in section_keywords:
        hits = page0.search_for(kw)
        if hits:
            start_y_pdf = hits[0].y0
            print(f"  섹션 시작 y: {start_y_pdf:.1f}pt ('{kw}')")
            break

    pages_needed = [pg_num]
    if pg_num + 1 < len(doc):
        next_text = doc[pg_num + 1].get_text()
        next_lines = next_text.split("\n")[:30]
        next_top = "\n".join(next_lines)
        if "지급시기" in next_top or "증권거래비용" in next_top or "합성총보수" in next_top:
            pages_needed.append(pg_num + 1)

    print(f"  캡처 페이지: {[p+1 for p in pages_needed]}")

    page_images = []
    page_heights_pdf = []
    for pn in pages_needed:
        p = doc[pn]
        pix = p.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
        page_images.append(Image.open(BytesIO(pix.tobytes("png"))))
        page_heights_pdf.append(p.rect.height)
    doc.close()

    if len(page_images) > 1:
        total_h = sum(img.height for img in page_images)
        full_img = Image.new("RGB", (page_images[0].width, total_h), "white")
        y_off = 0
        for img in page_images:
            full_img.paste(img, (0, y_off))
            y_off += img.height
    else:
        full_img = page_images[0]

    img_w, img_h = full_img.width, full_img.height
    single_page_img_h = page_images[0].height
    single_page_img_w = page_images[0].width

    try:
        with pdfplumber.open(pdf_path) as pdf:
            pl_page0 = pdf.pages[pg_num]
            pw0 = pl_page0.width
            ph0 = pl_page0.height
            sx = single_page_img_w / pw0
            sy = single_page_img_h / ph0

            x0_pdf, x1_pdf = None, None
            end_y_img = None
            fee_kws_cap = ["총보수", "총 보수", "증권거래비용", "증권 거래비용",
                           "거래비용", "합성총보수", "피투자"]
            num_re_cap  = re.compile(r"\d{1,2}\.\d{2,4}")

            for i, pn in enumerate(pages_needed):
                pl_page = pdf.pages[pn]
                y_offset_img = sum(page_images[j].height for j in range(i))

                for tbl in pl_page.find_tables():
                    ext = tbl.extract()
                    if not ext: continue
                    flat = " ".join(str(c) for row in ext for c in row if c)
                    if not any(kw in flat for kw in fee_kws_cap):
                        continue
                    total_valid = sum(
                        1 for row in ext for cell in row
                        for m in num_re_cap.findall(str(cell or ""))
                        if 0.0001 <= float(m) <= 9.9999
                    )
                    if total_valid < 2:
                        continue
                    bbox = tbl.bbox
                    if x0_pdf is None:
                        x0_pdf = bbox[0]
                        x1_pdf = bbox[2]
                    end_y_img = y_offset_img + int(bbox[3] * sy) + 20

            if x0_pdf is None:
                x0_pdf, x1_pdf = 50, pw0 - 50

            cx0 = max(0, int(x0_pdf * sx) - 10)
            cx1 = min(img_w, int(x1_pdf * sx) + 10)
            cy0 = max(0, int((start_y_pdf - 20) * sy)) if start_y_pdf else 0
            cy1 = end_y_img if end_y_img else img_h
            if cy1 <= cy0:   # 간이설명서 섹션 등 역전 방지
                cy1 = img_h

            print(f"  크롭 영역: ({cx0},{cy0}) → ({cx1},{cy1})")
            cropped = full_img.crop((cx0, cy0, cx1, cy1))
            print(f"  ✅ 크롭: {cropped.width}×{cropped.height}px")
    except Exception as e:
        print(f"  크롭 오류: {e}")
        cropped = full_img

    cropped.save(str(out_img), "PNG")
    print(f"  ✅ 이미지 저장: {out_img.name}")
    return out_img


def extract_fees(pdf_path: Path, pg_num: int) -> dict:
    result = {"synthetic_fee": None, "trading_cost": None}
    num_find = re.compile(r"\d{1,2}\.\d{2,4}")

    def is_valid_fee(s):
        try: return 0.0001 <= float(s) <= 9.9999
        except: return False

    pages_to_check = [pg_num, pg_num + 1]
    FEE_KWS = ["총보수", "총 보수", "증권거래비용", "증권 거래비용", "거래비용", "합성총보수"]

    with pdfplumber.open(pdf_path) as pdf:
        # Method 1: Table
        for chk_pg in pages_to_check:
            if chk_pg >= len(pdf.pages): continue
            pl_page = pdf.pages[chk_pg]
            for tbl in pl_page.find_tables():
                ext = tbl.extract()
                if not ext: continue
                flat = " ".join(str(c) for row in ext for c in row if c)
                if not any(kw in flat for kw in FEE_KWS): continue

                print(f"    표 발견 (p.{chk_pg+1}), 행 수: {len(ext)}")
                max_cols = max(len(row) for row in ext)

                if max_cols <= 5:  # 3열 형식 (구분 | 비율 | 시기)
                    for row in ext:
                        if not row or len(row) < 2: continue
                        key = str(row[0] or "").replace("\n", "").strip()
                        val_cells = [str(c or "") for c in row[1:] if c]
                        nums = [n for cell in val_cells for n in num_find.findall(cell) if is_valid_fee(n)]
                        if not nums: continue
                        if any(k in key for k in ["총보수·비용", "총보수·비", "합성총보수", "총 보수·비용",
                                                   "총 보수·비", "보수·비용"]):
                            if result["synthetic_fee"] is None:
                                result["synthetic_fee"] = nums[0]
                                print(f"    3열synth: '{key}'={nums[0]}")
                        if any(k in key for k in ["증권거래비용", "증권 거래비용", "거래비용"]):
                            if result["trading_cost"] is None:
                                result["trading_cost"] = nums[0]
                                print(f"    3열trade: '{key}'={nums[0]}")
                    if result["synthetic_fee"] and result["trading_cost"]:
                        return result
                else:  # wide 형식 (다열)
                    col_synth, col_trade = None, None
                    # "총보수·비용" 계열 (중간점 U+00B7 및 전각 U+FF65 모두 처리)
                    SYNTH_KWS = ["피투자", "보수포함",
                                 "총보수·비용", "총보수·비", "총보수·",
                                 "총보수･비용", "총보수･비", "총보수･"]
                    # 헤더가 여러 행에 나뉘어 있는 경우를 위해 열별 텍스트를 concat
                    col_concat: dict[int, str] = {}
                    num_re_hdr = re.compile(r"^[\d.,%-]+$")
                    for row in ext:
                        if not row: continue
                        for ci, cell in enumerate(row):
                            t = str(cell or "").replace("\n", "").strip()
                            # 순수 숫자 셀은 헤더 concat 에서 제외
                            if t and not num_re_hdr.match(t):
                                col_concat[ci] = col_concat.get(ci, "") + t
                    for ci, t in col_concat.items():
                        if col_synth is None and any(k in t for k in SYNTH_KWS):
                            col_synth = ci
                        if col_trade is None and any(k in t for k in ["증권거래비용", "거래비용"]):
                            col_trade = ci

                    print(f"    열 인덱스: col_synth={col_synth}, col_trade={col_trade}")

                    flat_nums = []
                    if col_synth is not None:
                        flat_nums = [n for row in ext for ci, cell in enumerate(row)
                                     if ci == col_synth for n in num_find.findall(str(cell or ""))
                                     if is_valid_fee(n)]
                        print(f"    데이터 행 숫자(유효): {flat_nums}")

                    if flat_nums:
                        if col_synth != col_trade and col_trade is not None:
                            trade_nums = [n for row in ext for ci, cell in enumerate(row)
                                          if ci == col_trade for n in num_find.findall(str(cell or ""))
                                          if is_valid_fee(n)]
                            if flat_nums and trade_nums and flat_nums[-1] != trade_nums[-1]:
                                result["synthetic_fee"] = flat_nums[-1]
                                result["trading_cost"]  = trade_nums[-1]
                                print(f"    ✅ 두열인덱스: synth={flat_nums[-1]}, trade={trade_nums[-1]}")
                                return result
                        result["synthetic_fee"] = flat_nums[-1]
                        print(f"    ✅ col_synth: synth={flat_nums[-1]}")
                    else:
                        # wide 테이블이지만 row[0]=레이블 형식인 경우 (ACE 등)
                        print(f"    wide col 데이터 없음 → row-based fallback")
                        ROW_SYNTH_KWS = ["총보수·비용", "총보수·비", "총보수･비용", "총보수･비",
                                         "합성총보수", "합성 총보수", "총 보수·비용", "보수·비용"]
                        ROW_TRADE_KWS = ["증권거래비용", "증권 거래비용", "거래비용"]
                        for row in ext:
                            if not row or len(row) < 2: continue
                            key = str(row[0] or "").replace("\n", "").strip()
                            val_cells = [str(c or "") for c in row[1:] if c]
                            nums = [n for cell in val_cells
                                    for n in num_find.findall(cell) if is_valid_fee(n)]
                            if not nums: continue
                            if result["synthetic_fee"] is None and any(k in key for k in ROW_SYNTH_KWS):
                                result["synthetic_fee"] = nums[0]
                                print(f"    wide-row synth: '{key}'={nums[0]}")
                            if result["trading_cost"] is None and any(k in key for k in ROW_TRADE_KWS):
                                result["trading_cost"] = nums[0]
                                print(f"    wide-row trade: '{key}'={nums[0]}")

        # Method 2: Text
        print("    텍스트 방식으로 fallback...")
        for chk_pg in pages_to_check:
            if chk_pg >= len(pdf.pages): continue
            text = pdf.pages[chk_pg].extract_text() or ""
            for line in text.split("\n"):
                line = line.strip()
                nums = [n for n in num_find.findall(line) if is_valid_fee(n)]
                if not nums: continue
                if any(k in line for k in ["합성총보수", "합성 총보수",
                                            "총보수·비용", "총보수·비", "총 보수·비용",
                                            "총보수･비용", "총보수･비",
                                            "피투자", "보수포함"]):
                    if result["synthetic_fee"] is None:
                        result["synthetic_fee"] = nums[-1]
                if any(k in line for k in ["증권거래비용", "증권 거래비용"]):
                    if result["trading_cost"] is None:
                        result["trading_cost"] = nums[-1]
                if len(nums) >= 8 and result["synthetic_fee"] is None:
                    result["synthetic_fee"] = nums[-2]
                    result["trading_cost"]  = nums[-1]
                    print(f"    숫자8+행: {nums}")

    return result


# ═══════════════════════════════════════════════════════════════════════════
# ── SECTION 2: 미국 ETF (etfcheck.co.kr) ────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════

async def scrape_us_etf(page, item: dict) -> dict:
    ticker     = item["ticker"]
    etf_id     = item.get("etfcheck_id")
    result     = {
        "ticker": ticker, "name": item["name"],
        "설정일": None, "시가총액": None, "TER": None,
        "img": None
    }

    if not etf_id:
        print(f"  ⚠️ {ticker}: etfcheck ID 없음")
        return result

    url = f"https://www.etfcheck.co.kr/mobile/global/etpitem/{etf_id}/basic/개요"
    print(f"  [etfcheck] {ticker}: {url}")
    await page.goto(url, timeout=25000)
    await page.wait_for_load_state("networkidle", timeout=20000)
    await asyncio.sleep(3)

    # 텍스트에서 데이터 추출
    text = await page.evaluate("() => document.body.innerText")

    m = re.search(r"설정일\s+(\d{4}\.\d{2}\.\d{2})", text)
    if m: result["설정일"] = m.group(1)

    m = re.search(r"시가총액\s+(\$\s*[\d,]+\.?\d*[A-Z]?)", text)
    if m: result["시가총액"] = m.group(1).strip()

    m = re.search(r"TER\s+([\d.]+%)", text)
    if m: result["TER"] = m.group(1)

    print(f"  ✅ {ticker}: 설정일={result['설정일']}, 시가총액={result['시가총액']}, TER={result['TER']}")

    # 스크린샷: ETF 개요(설정일/시가총액/TER) 섹션이 보이도록 스크롤
    img_path = IMG_DIR / f"{ticker}_etfcheck.png"
    try:
        # "ETF 개요" 또는 "설정일" 텍스트가 있는 요소로 스크롤
        scrolled = await page.evaluate("""() => {
            const els = document.querySelectorAll('*');
            for (const el of els) {
                if (el.children.length === 0 &&
                    (el.textContent.trim() === 'ETF 개요' ||
                     el.textContent.trim() === '설정일')) {
                    el.scrollIntoView({behavior: 'instant', block: 'start'});
                    return true;
                }
            }
            window.scrollBy(0, 480);
            return false;
        }""")
        await asyncio.sleep(1)
    except Exception:
        pass
    await page.screenshot(path=str(img_path), full_page=False)
    result["img"] = img_path
    print(f"  ✅ 스크린샷 저장: {img_path.name}")

    return result


# ═══════════════════════════════════════════════════════════════════════════
# ── SECTION 3: PPT 생성 ──────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════

TITLE_COLOR  = RGBColor(0x1F, 0x49, 0x7D)
HEADER_COLOR = RGBColor(0x1F, 0x49, 0x7D)
ROW_ALT      = RGBColor(0xE9, 0xF0, 0xF8)
ROW_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
FONT_KR      = "맑은 고딕"


def set_cell_text(cell, text, font_size=11, bold=False,
                  bg_color=None, align=PP_ALIGN.CENTER, color=None):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = FONT_KR
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color

    if bg_color:
        from lxml import etree
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        solidFill = etree.SubElement(tcPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
        srgbClr = etree.SubElement(solidFill, "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
        srgbClr.set("val", str(bg_color))


def add_title_slide_kr(prs, kr_results):
    """슬라이드 1: 한국 ETF 보수 요약"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width
    sh = prs.slide_height

    # 제목
    txb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), sw - Inches(0.8), Inches(0.6))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "한국 ETF 보수 요약 (우주항공)"
    run.font.name = FONT_KR
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR

    # 표 (헤더 + 행)
    headers = ["상품명", "티커", "총보수·비용\n(피투자포함)", "증권거래비용"]
    col_w = [Inches(3.5), Inches(1.1), Inches(1.8), Inches(1.8)]
    t_left = Inches(0.4)
    t_top  = Inches(1.0)
    n_rows = 1 + len(kr_results)
    tbl = slide.shapes.add_table(n_rows, 4, t_left, t_top,
                                  sum(col_w), Inches(0.38 * n_rows)).table

    for ci, w in enumerate(col_w):
        tbl.columns[ci].width = w

    # 헤더
    for ci, h in enumerate(headers):
        set_cell_text(tbl.cell(0, ci), h, font_size=11, bold=True,
                      bg_color=HEADER_COLOR, color=RGBColor(0xFF, 0xFF, 0xFF))

    # 데이터 행
    for ri, r in enumerate(kr_results, 1):
        bg = ROW_ALT if ri % 2 == 0 else ROW_WHITE
        sf = f"{r['synthetic_fee']}%" if r.get("synthetic_fee") else "N/A"
        tc = f"{r['trading_cost']}%"  if r.get("trading_cost")  else "N/A"
        set_cell_text(tbl.cell(ri, 0), r["name"],  font_size=11, bg_color=bg, align=PP_ALIGN.LEFT)
        set_cell_text(tbl.cell(ri, 1), r["sotCd"], font_size=11, bg_color=bg)
        set_cell_text(tbl.cell(ri, 2), sf,          font_size=11, bg_color=bg)
        set_cell_text(tbl.cell(ri, 3), tc,          font_size=11, bg_color=bg)

    return slide


def add_title_slide_us(prs, us_results):
    """슬라이드 2: 미국 ETF 기본정보 요약"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width
    sh = prs.slide_height

    txb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), sw - Inches(0.8), Inches(0.6))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "미국 ETF 기본정보 요약 (우주항공)"
    run.font.name = FONT_KR
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR

    headers = ["상품명", "티커", "설정일", "시가총액", "TER"]
    col_w = [Inches(3.0), Inches(0.8), Inches(1.2), Inches(1.8), Inches(1.0)]
    t_left = Inches(0.4)
    t_top  = Inches(1.0)
    n_rows = 1 + len(us_results)
    tbl = slide.shapes.add_table(n_rows, 5, t_left, t_top,
                                  sum(col_w), Inches(0.38 * n_rows)).table

    for ci, w in enumerate(col_w):
        tbl.columns[ci].width = w

    for ci, h in enumerate(headers):
        set_cell_text(tbl.cell(0, ci), h, font_size=11, bold=True,
                      bg_color=HEADER_COLOR, color=RGBColor(0xFF, 0xFF, 0xFF))

    for ri, r in enumerate(us_results, 1):
        bg = ROW_ALT if ri % 2 == 0 else ROW_WHITE
        set_cell_text(tbl.cell(ri, 0), r["name"],       font_size=10, bg_color=bg, align=PP_ALIGN.LEFT)
        set_cell_text(tbl.cell(ri, 1), r["ticker"],     font_size=11, bg_color=bg)
        set_cell_text(tbl.cell(ri, 2), r["설정일"] or "N/A",  font_size=10, bg_color=bg)
        set_cell_text(tbl.cell(ri, 3), r["시가총액"] or "N/A", font_size=10, bg_color=bg)
        set_cell_text(tbl.cell(ri, 4), r["TER"]  or "N/A",   font_size=11, bg_color=bg)

    return slide


def add_kr_etf_slide(prs, item, fees, img_path):
    """한국 ETF 개별 슬라이드: 보수 표 이미지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width
    sh = prs.slide_height

    # 제목
    txb = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), sw - Inches(0.6), Inches(0.45))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = item["name"]
    run.font.name = FONT_KR
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR

    # 보수 요약
    sf = f"{fees.get('synthetic_fee', 'N/A')}%" if fees.get("synthetic_fee") else "N/A"
    tc = f"{fees.get('trading_cost',  'N/A')}%" if fees.get("trading_cost")  else "N/A"
    subtitle = f"총보수·비용(피투자포함): {sf}    증권거래비용: {tc}"
    txb2 = slide.shapes.add_textbox(Inches(0.3), Inches(0.62), sw - Inches(0.6), Inches(0.35))
    tf2 = txb2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.name = FONT_KR
    run2.font.size = Pt(11)
    run2.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # 이미지
    if img_path and Path(img_path).exists():
        img = Image.open(img_path)
        iw, ih = img.size
        avail_w = sw - Inches(0.6)
        avail_h = sh - Inches(1.15)
        scale = min(avail_w / (iw * 9144), avail_h / (ih * 9144))
        disp_w = int(iw * 9144 * scale)
        disp_h = int(ih * 9144 * scale)
        left = (sw - disp_w) // 2
        slide.shapes.add_picture(str(img_path), left, Inches(1.05), disp_w, disp_h)

    return slide


def add_us_etf_slide(prs, us_result):
    """미국 ETF 개별 슬라이드: etfcheck 스크린샷"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width
    sh = prs.slide_height

    ticker = us_result["ticker"]
    name   = us_result["name"]

    # 제목
    txb = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), sw - Inches(0.6), Inches(0.45))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{ticker}  ({name})"
    run.font.name = FONT_KR
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR

    # 기본정보 요약
    info = f"설정일: {us_result['설정일'] or 'N/A'}    시가총액: {us_result['시가총액'] or 'N/A'}    TER: {us_result['TER'] or 'N/A'}"
    txb2 = slide.shapes.add_textbox(Inches(0.3), Inches(0.62), sw - Inches(0.6), Inches(0.35))
    tf2 = txb2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = info
    run2.font.name = FONT_KR
    run2.font.size = Pt(11)
    run2.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # 스크린샷 이미지
    img_path = us_result.get("img")
    if img_path and Path(img_path).exists():
        img = Image.open(img_path)
        iw, ih = img.size
        avail_w = sw - Inches(0.6)
        avail_h = sh - Inches(1.15)
        scale = min(avail_w / (iw * 9144), avail_h / (ih * 9144))
        disp_w = int(iw * 9144 * scale)
        disp_h = int(ih * 9144 * scale)
        left = (sw - disp_w) // 2
        slide.shapes.add_picture(str(img_path), left, Inches(1.05), disp_w, disp_h)
    else:
        txb3 = slide.shapes.add_textbox(Inches(0.3), Inches(1.3), sw - Inches(0.6), Inches(2))
        tf3 = txb3.text_frame
        tf3.paragraphs[0].add_run().text = "etfcheck.co.kr 데이터 없음 (N/A)"

    return slide


def build_ppt(kr_items_with_fees, us_results):
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1: 한국 ETF 요약
    add_title_slide_kr(prs, kr_items_with_fees)

    # Slide 2: 미국 ETF 요약
    add_title_slide_us(prs, us_results)

    # Slides 3-6: 한국 ETF 개별
    for item_fees in kr_items_with_fees:
        add_kr_etf_slide(prs, item_fees, item_fees, item_fees.get("img"))

    # Slides 7-10: 미국 ETF 개별
    for r in us_results:
        add_us_etf_slide(prs, r)

    out = OUT_DIR / f"space_etf_{TIMESTAMP}.pptx"
    prs.save(str(out))
    print(f"\n✅ PPT 저장: {out}")
    return out


# ═══════════════════════════════════════════════════════════════════════════
# ── MAIN ────────────────────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════

async def main():
    print("=" * 60)
    print("우주항공 ETF 보고서 생성기")
    print("=" * 60)

    kr_items_with_fees = []
    us_results = []

    async with async_playwright() as p:
        # ── 브라우저 1: funetf (한국 ETF) ──────────────────────────────
        browser_kr = await p.chromium.launch(headless=True)
        context_kr = await browser_kr.new_context(
            user_agent='Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 Chrome/120.0.0.0 Safari/537.36',
            accept_downloads=True,
            viewport={"width": 1400, "height": 900},
        )
        page_kr = await context_kr.new_page()

        print("\n[로그인]")
        await login(page_kr)

        for item in KR_ITEMS:
            name   = item["name"]
            sotCd  = item["sotCd"]
            safe   = name.replace(" ", "_").replace("/", "_")
            pdf_path = PDF_DIR / f"{safe}.pdf"
            img_path = IMG_DIR / f"{safe}_fee.png"

            print(f"\n{'='*50}")
            print(f"▶ {name} ({sotCd})")
            print(f"{'='*50}")

            # PDF 다운로드 (이미 있으면 재사용)
            if pdf_path.exists() and pdf_path.stat().st_size > 10000:
                print(f"  PDF 재사용: {pdf_path.name}")
            else:
                ok = await download_kr_pdf(page_kr, context_kr, sotCd, pdf_path)
                if not ok:
                    kr_items_with_fees.append({**item, "synthetic_fee": None, "trading_cost": None, "img": None})
                    continue

            # 보수 페이지 찾기 + 캡처
            print("  [보수 페이지 탐색]")
            pg_num = find_fee_page(pdf_path)
            if pg_num is None:
                kr_items_with_fees.append({**item, "synthetic_fee": None, "trading_cost": None, "img": None})
                continue

            capture_fee_table(pdf_path, pg_num, img_path)
            fees = extract_fees(pdf_path, pg_num)

            print(f"  ▶ 합성총보수:   {fees['synthetic_fee']}%")
            print(f"  ▶ 증권거래비용: {fees['trading_cost']}%")

            kr_items_with_fees.append({
                **item,
                "synthetic_fee": fees["synthetic_fee"],
                "trading_cost":  fees["trading_cost"],
                "img": str(img_path) if img_path.exists() else None,
            })

        await browser_kr.close()

        # ── 브라우저 2: etfcheck (미국 ETF) ────────────────────────────
        browser_us = await p.chromium.launch(headless=True)
        context_us = await browser_us.new_context(
            user_agent='Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 Chrome/120.0.0.0 Safari/537.36',
            viewport={"width": 1200, "height": 900},
        )
        page_us = await context_us.new_page()

        print(f"\n\n{'='*50}")
        print("▶ 미국 ETF (etfcheck.co.kr)")
        print(f"{'='*50}")

        for item in US_ITEMS:
            print(f"\n  [{item['ticker']}]")
            result = await scrape_us_etf(page_us, item)
            us_results.append(result)

        await browser_us.close()

    # ── PPT 생성 ──────────────────────────────────────────────────────
    print("\n\n[PPT 생성]")
    ppt_path = build_ppt(kr_items_with_fees, us_results)

    # ── 최종 요약 ─────────────────────────────────────────────────────
    print("\n" + "=" * 60)
    print("한국 ETF 보수 요약")
    print("-" * 60)
    for r in kr_items_with_fees:
        sf = f"{r['synthetic_fee']}%" if r.get("synthetic_fee") else "N/A"
        tc = f"{r['trading_cost']}%"  if r.get("trading_cost")  else "N/A"
        print(f"  {r['name']:30s} {sf:10s} {tc}")

    print("\n" + "=" * 60)
    print("미국 ETF 기본정보 요약")
    print("-" * 60)
    for r in us_results:
        print(f"  {r['ticker']:6s} 설정일={r['설정일'] or 'N/A':12s} 시가총액={r['시가총액'] or 'N/A':15s} TER={r['TER'] or 'N/A'}")


asyncio.run(main())
