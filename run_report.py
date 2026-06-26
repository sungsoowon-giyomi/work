"""
ETF 보고서 런처
config.txt 를 읽어 space_etf_report.py 를 실행합니다.
"""
import asyncio
import json
import re
import requests
from io import BytesIO
from pathlib import Path
from datetime import datetime

import fitz
import pdfplumber
from PIL import Image, ImageDraw
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 인증 ──────────────────────────────────────────────────────────────────
EMAIL    = "sungsoowon45@gmail.com"
PASSWORD = "!sungsoo0405"
ETFCHECK_EMAIL = "sswon9545@gmail.com"
ETFCHECK_PW    = "tjdtn0405"

OUT_DIR  = Path("output/space_etf")
PDF_DIR  = OUT_DIR / "pdfs"
IMG_DIR  = OUT_DIR / "imgs"
OUT_DIR.mkdir(parents=True, exist_ok=True)
PDF_DIR.mkdir(exist_ok=True)
IMG_DIR.mkdir(exist_ok=True)

TIMESTAMP = datetime.now().strftime("%Y%m%d_%H%M")


# ═══════════════════════════════════════════════════════════════════════════
# ── config.txt 파싱 ──────────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════

def load_config(path: str = "config.txt"):
    kr_items = []
    us_items = []
    fund_items = []
    section = None

    with open(path, encoding="utf-8") as f:
        for raw in f:
            line = raw.strip()
            if not line or line.startswith("#"):
                continue
            if line == "[한국ETF]":
                section = "kr"
                continue
            if line == "[미국ETF]":
                section = "us"
                continue
            if line == "[펀드]":
                section = "fund"
                continue

            if section == "kr":
                parts = line.split(None, 1)
                code = parts[0].lstrip("Aa")   # 'A' 접두어 제거
                name = parts[1].strip() if len(parts) > 1 else code
                kr_items.append({"name": name, "sotCd": code})

            elif section == "us":
                parts = line.split()
                ticker = parts[0].upper()
                etfcheck_id = None
                name = ticker
                if len(parts) >= 2:
                    if parts[1].upper() == "N/A":
                        etfcheck_id = "N/A"
                    else:
                        name = parts[1]
                if len(parts) >= 3 and parts[2].startswith("F"):
                    etfcheck_id = parts[2]
                us_items.append({"ticker": ticker, "name": name, "_preset_id": etfcheck_id})

            elif section == "fund":
                # 형식: FUND_ID [CLASS] [펀드명]
                # 예: K55105BA7311
                #     K55105BA7311 A
                #     K55105BA7311 C-Pe 삼성미국S&P500펀드
                tokens = line.split()
                fund_id = None
                fund_class = None
                name_parts = []
                for tok in tokens:
                    if re.match(r'^K[0-9A-Z]{10,}$', tok, re.IGNORECASE):
                        fund_id = tok.upper()
                    elif fund_id and re.match(r'^[A-Za-z][A-Za-z0-9\-]{0,6}$', tok) and fund_class is None:
                        fund_class = tok          # 펀드코드 뒤 짧은 영문 토큰 = 클래스명
                    else:
                        name_parts.append(tok)
                fname = ' '.join(name_parts) or fund_id or line
                fund_items.append({"name": fname, "sotCd": None, "_fund_id": fund_id, "_class": fund_class})

    print(f"\n[config] 한국 ETF {len(kr_items)}개, 미국 ETF {len(us_items)}개, 펀드 {len(fund_items)}개 로드")
    return kr_items, us_items, fund_items


# ═══════════════════════════════════════════════════════════════════════════
# ── 미국 ETF: etfcheck ID 자동 검색 ─────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════

async def find_etfcheck_id(page, ticker: str, _retry: int = 0) -> tuple[str, str] | None:
    """
    etfcheck.co.kr 검색 → 첫 번째 해외 ETF 결과 클릭 → URL에서 펀드 ID 추출
    """
    await page.goto("https://www.etfcheck.co.kr/mobile/main", timeout=25000)
    await page.wait_for_load_state("networkidle", timeout=20000)
    await asyncio.sleep(2)

    # Vuetify 오버레이 제거
    await page.evaluate("document.querySelectorAll('.v-overlay').forEach(e=>e.remove())")
    await asyncio.sleep(0.3)

    try:
        # 검색 아이콘 버튼 클릭 → 검색 모달 열기
        await page.locator('button[class*="topicon_search"]').click(force=True, timeout=5000)
        await asyncio.sleep(1)

        # 모달 안 editable input에 티커 입력
        editable = page.locator('input[type=text]:not([readonly])').first
        await editable.fill(ticker, timeout=5000)
        await page.keyboard.press("Enter")
        await asyncio.sleep(3)
        await page.wait_for_load_state("networkidle", timeout=10000)
        await asyncio.sleep(2)

        # 검색 결과 페이지 텍스트에서 ETF명 추출
        # 화면 형태: "NASA 월\nTema Space Innovators ETF  31.60  ..."
        body_text = await page.evaluate("() => document.body.innerText")
        lines = [l.strip() for l in body_text.split("\n") if l.strip()]
        etf_name = None
        for i, line in enumerate(lines):
            if line == ticker or line.startswith(ticker + " "):
                for j in range(i + 1, min(i + 3, len(lines))):
                    cand = lines[j]
                    if "ETF" in cand or "Trust" in cand or "Fund" in cand:
                        etf_name = re.split(r"\s{2,}|\t", cand)[0].strip()
                        break
                if etf_name:
                    break

        if not etf_name:
            print(f"    ⚠️  {ticker}: 검색 결과에서 ETF명을 찾지 못함")
            return None

        await page.get_by_text(etf_name, exact=False).first.click(timeout=5000)
        await asyncio.sleep(2)

    except Exception as e:
        if _retry < 1:
            print(f"    검색 오류: {e} → 재시도")
            await asyncio.sleep(2)
            return await find_etfcheck_id(page, ticker, _retry + 1)
        print(f"    검색 오류: {e}")
        return None

    # 이동된 URL에서 펀드 ID 추출
    current_url = page.url
    m = re.search(r"etpitem/([^/]+)", current_url)
    if m:
        fund_id = m.group(1)
        print(f"    ✅ {ticker} → {fund_id} ({etf_name})")
        return fund_id, etf_name

    print(f"    ⚠️  {ticker}: 자동 검색 실패 → config.txt에 ID 직접 지정 가능")
    print(f"       (현재 URL: {current_url})")
    return None


# ═══════════════════════════════════════════════════════════════════════════
# ── 한국 ETF 관련 함수 (space_etf_report.py 와 동일) ─────────────────────
# ═══════════════════════════════════════════════════════════════════════════

def summarize(text: str, limit: int = 50) -> str:
    text = re.sub(r"\s+", " ", text or "").strip()
    if len(text) <= limit:
        return text
    return text[:limit - 1].rstrip() + "…"


async def login(page):
    await page.goto("https://www.funetf.co.kr/auth")
    await page.wait_for_load_state("domcontentloaded")
    await page.fill('input[name="username"]', EMAIL)
    await page.fill('input[name="password"]', PASSWORD)
    await page.locator('button:has-text("로그인")').last.click()
    await page.wait_for_url("https://www.funetf.co.kr/", timeout=15000)
    print("  ✅ funetf 로그인 성공")


async def etfcheck_login(page):
    """etfcheck 로그인 (Vuetify: click → type → JS button click)"""
    await page.goto("https://www.etfcheck.co.kr/mobile/user/signin", timeout=20000)
    await page.wait_for_load_state("networkidle", timeout=15000)
    await asyncio.sleep(2)

    await page.locator("input[type=email]").click(timeout=5000)
    await page.keyboard.type(ETFCHECK_EMAIL)
    await page.locator("input[type=password]").click(timeout=5000)
    await page.keyboard.type(ETFCHECK_PW)
    await asyncio.sleep(1)

    await page.evaluate("""
        (() => {
            const btns = Array.from(document.querySelectorAll('button'));
            const btn = btns.find(b => b.textContent.trim() === '로그인');
            if (btn) btn.click();
        })()
    """)
    await asyncio.sleep(5)
    await page.wait_for_load_state("networkidle", timeout=20000)
    print("  ✅ etfcheck 로그인 성공")


async def find_kr_product_link(page, sotCd: str, name: str | None = None) -> str | None:
    """funetf에서 한국 ETF 상품 링크 찾기.
    코드(숫자 또는 알파벳 혼합 모두)로 먼저 검색, 실패하면 이름으로 폴백.
    """
    from urllib.parse import quote

    async def _search(query: str) -> str | None:
        search_url = f"https://www.funetf.co.kr/search?schVal={quote(query)}"
        await page.goto(search_url)
        await page.wait_for_load_state("networkidle", timeout=25000)
        await asyncio.sleep(2)
        return await page.evaluate('''() => {
            for (const a of document.querySelectorAll('a')) {
                if (a.href.includes('/product/etf/view/') || a.href.includes('/view/etf/')) return a.href;
            }
            return null;
        }''')

    # 코드(숫자만 or 알파벳 혼합 모두)로 먼저 검색
    link = await _search(sotCd)
    if link:
        return link

    # 코드 검색 실패 시 이름으로 폴백
    if name and name != sotCd:
        return await _search(name)

    return None


async def find_fund_link(page, name: str, fund_id: str | None = None) -> str | None:
    """funetf 일반펀드 상품 링크 찾기.
    fund_id가 있으면 직접 이동, 없으면 이름으로 검색.
    """
    if fund_id:
        url = f"https://www.funetf.co.kr/product/fund/view/{fund_id}"
        await page.goto(url, timeout=20000)
        await page.wait_for_load_state("networkidle", timeout=20000)
        await asyncio.sleep(2)
        return url if "/fund/view/" in page.url else None

    from urllib.parse import quote
    search_url = f"https://www.funetf.co.kr/search?schVal={quote(name)}"
    await page.goto(search_url)
    await page.wait_for_load_state("networkidle", timeout=25000)
    await asyncio.sleep(2)

    return await page.evaluate('''() => {
        for (const a of document.querySelectorAll('a')) {
            if (a.href && a.href.includes('/product/fund/view/')) return a.href;
        }
        return null;
    }''')


async def scrape_fund_detail(page, name: str, fund_id: str | None = None) -> dict:
    """funetf 일반펀드 상세페이지에서 운용사/투자위험/설정일/순자산/특징 추출"""
    result = {"펀드명": None, "운용사": None, "투자위험": None, "상장일": None, "순자산": None, "특징": None, "holdings": []}

    link = await find_fund_link(page, name, fund_id)
    if not link:
        print(f"  ❌ '{name}' 상품 링크 없음 (검색 결과 없음)")
        return result

    await page.goto(link)
    await page.wait_for_load_state("networkidle", timeout=20000)
    await asyncio.sleep(2)

    text = await page.evaluate("() => document.body.innerText")

    # 실제 펀드명 추출 (페이지 h1 또는 상품명 영역)
    fund_name = await page.evaluate("""() => {
        const h1 = document.querySelector('h1, .product-name, .prd-name, [class*="title"]');
        return h1 ? h1.textContent.trim() : null;
    }""")
    if fund_name and len(fund_name) > 5:
        result["펀드명"] = fund_name

    m = re.search(r"운용사\n(.+)", text)
    if m: result["운용사"] = m.group(1).strip()

    # 투자위험: "3등급" 다음 줄에 "(다소높은위험)" 있을 수 있음
    m = re.search(r"투자위험\n(\d+등급)", text)
    if m:
        risk = m.group(1).strip()
        m2 = re.search(r"투자위험\n\d+등급\n(\([가-힣]+위험\))", text)
        result["투자위험"] = risk + (" " + m2.group(1) if m2 else "")

    # 순자산: "순자산 기준 : 858 억원" 또는 "순자산\n..." 형식
    m = re.search(r"순자산 기준 : (.+)", text)
    if m:
        result["순자산"] = m.group(1).strip()
    else:
        m = re.search(r"순자산\n(.+)", text)
        if m: result["순자산"] = m.group(1).strip()

    # 설정일 또는 상장일
    m = re.search(r"설정일\n(\d{4}\.\d{2}\.\d{2})", text)
    if m:
        result["상장일"] = m.group(1)
    else:
        m = re.search(r"상장일\n(\d{4}\.\d{2}\.\d{2})", text)
        if m: result["상장일"] = m.group(1)

    try:
        await page.locator("#btnInfoDetail").click(timeout=5000)
        await asyncio.sleep(1.5)
        feature_text = await page.evaluate("""() => {
            const parts = [];
            for (const id of ['divDiscription1', 'divDiscription2']) {
                const el = document.getElementById(id);
                if (el) parts.push(el.innerText);
            }
            return parts.join(' ');
        }""")
        if feature_text:
            feature_text = feature_text.replace("펀드특징", "").replace("운용전략", "")
            result["특징"] = summarize(feature_text, 50)
    except Exception:
        pass

    return result


async def download_fund_pdf(page, context, name: str, save_path: Path, fund_id: str | None = None) -> bool:
    """funetf 일반펀드 투자설명서 PDF 다운로드"""
    link = await find_fund_link(page, name, fund_id)
    if not link:
        return False

    await page.goto(link)
    await page.wait_for_load_state("networkidle", timeout=20000)
    await asyncio.sleep(3)

    dl_selector = 'div.prd-11-dcmt-item:has(p.prd-11-dcmt-title:text("투자설명서")) button:has-text("다운로드")'
    try:
        btn = page.locator(dl_selector)
        if await btn.count() > 0:
            async with page.expect_download(timeout=30000) as dl_info:
                await btn.first.click()
            dl = await dl_info.value
            await dl.save_as(str(save_path))
            print(f"  ✅ PDF 저장: {save_path.name}")
            return True
    except Exception as e:
        print(f"  버튼 클릭 오류: {e}")

    try:
        fileurl = await page.evaluate('''() => {
            const items = document.querySelectorAll(".prd-11-dcmt-item");
            for (const item of items) {
                const title = item.querySelector(".prd-11-dcmt-title");
                if (title && title.textContent?.trim() === "투자설명서") {
                    const btn = item.querySelector("button.download");
                    return btn ? btn.getAttribute("data-fileurl") : null;
                }
            }
            return null;
        }''')
        if fileurl:
            full_url = f"https://www.funetf.co.kr{fileurl}"
            resp = requests.get(full_url, timeout=30)
            if resp.ok and len(resp.content) > 1000:
                save_path.write_bytes(resp.content)
                print(f"  ✅ 저장(fallback): {save_path.name}")
                return True
    except Exception as e:
        print(f"  fallback 오류: {e}")

    return False


async def capture_fund_fee_screenshot(page, fund_url: str, img_path: Path) -> bool:
    """funetf 펀드 보수 탭의 보수 테이블 전체를 스크린샷으로 저장"""
    await page.goto(fund_url, timeout=20000)
    await page.wait_for_load_state("networkidle", timeout=20000)
    await asyncio.sleep(2)

    # 보수 탭 클릭 — 탭/메뉴 항목 중 "보수" 관련 텍스트를 포함하는 요소 클릭
    clicked = await page.evaluate("""() => {
        // 탭 영역에 해당할 가능성이 높은 요소 우선 탐색
        const tabSelectors = [
            '.tab-item', '.tab_item', '.nav-item', '.menu-item',
            'ul.tab li', 'ul.tabs li', 'ul.nav li',
            '.prd-tab li', '.product-tab li',
        ];
        const keywords = ['보수', '수수료', '비용'];
        for (const sel of tabSelectors) {
            for (const el of document.querySelectorAll(sel)) {
                const txt = el.textContent?.trim() || '';
                if (keywords.some(k => txt.includes(k))) {
                    // li 안의 a/button 클릭 시도
                    const inner = el.querySelector('a, button') || el;
                    inner.click();
                    return txt;
                }
            }
        }
        // 폴백: 모든 a/button/li/span에서 텍스트 매칭
        for (const el of document.querySelectorAll('a, button, li, span')) {
            const txt = el.textContent?.trim() || '';
            if (keywords.some(k => txt === k || txt === k + '/수수료')) {
                el.click();
                return txt;
            }
        }
        return null;
    }""")
    if clicked:
        print(f"  ▶ '{clicked}' 탭 클릭")
        await asyncio.sleep(2)
        await page.wait_for_load_state("networkidle", timeout=10000)
        await asyncio.sleep(1)
    else:
        # 탭 클릭 실패 시 페이지 내 탭 목록 출력해서 디버그
        tabs = await page.evaluate("""() => {
            const items = [];
            for (const el of document.querySelectorAll('.tab-item, .tab_item, .nav-item, ul.tab li, ul.tabs li')) {
                items.push(el.textContent?.trim());
            }
            return items.slice(0, 10);
        }""")
        print(f"  ⚠️ 보수 탭 못 찾음. 발견된 탭: {tabs}")
        print(f"  → 현재 페이지에서 보수 테이블 탐색 계속")

    # 보수 테이블 요소 찾기
    table_box = await page.evaluate("""() => {
        // 테이블 포함 영역 탐색
        const tables = document.querySelectorAll('table');
        for (const t of tables) {
            const txt = t.innerText || '';
            if (txt.includes('총보수') || txt.includes('보수') || txt.includes('클래스')) {
                const rect = t.getBoundingClientRect();
                // 테이블 위 섹션 헤더까지 포함하도록 여백 추가
                return {x: rect.left, y: Math.max(0, rect.top - 60),
                        width: rect.width, height: rect.height + 80};
            }
        }
        // 테이블 없으면 보수 관련 섹션 전체
        const sections = document.querySelectorAll('section, .section, .tab-content, [class*="fee"], [class*="cost"]');
        for (const s of sections) {
            const txt = s.innerText || '';
            if (txt.includes('총보수') || txt.includes('보수')) {
                const rect = s.getBoundingClientRect();
                return {x: rect.left, y: Math.max(0, rect.top - 20),
                        width: rect.width, height: rect.height + 40};
            }
        }
        return null;
    }""")

    # 탭 클릭 없이도 페이지 내 "보수" 섹션으로 스크롤 시도
    if not table_box:
        await page.evaluate("""() => {
            for (const el of document.querySelectorAll('h2, h3, h4, th, td, p, span, div')) {
                const txt = el.textContent?.trim() || '';
                if (txt === '보수' || txt.includes('총보수') || txt.includes('보수율')) {
                    el.scrollIntoView({behavior: 'instant', block: 'start'});
                    return;
                }
            }
        }""")
        await asyncio.sleep(0.5)
        table_box = await page.evaluate("""() => {
            const tables = document.querySelectorAll('table');
            for (const t of tables) {
                const txt = t.innerText || '';
                if (txt.includes('총보수') || txt.includes('보수') || txt.includes('클래스')) {
                    const rect = t.getBoundingClientRect();
                    return {x: rect.left, y: Math.max(0, rect.top - 60),
                            width: rect.width, height: rect.height + 80};
                }
            }
            return null;
        }""")

    if table_box and table_box.get("height", 0) > 50:
        # 해당 요소가 뷰포트에 들어오도록 스크롤
        await page.evaluate(f"window.scrollTo(0, {max(0, table_box['y'] - 40)})")
        await asyncio.sleep(0.5)
        # clip 으로 테이블만 캡처
        clip = {
            "x": max(0, table_box["x"] - 10),
            "y": max(0, table_box["y"] - 40),
            "width": min(table_box["width"] + 20, 1920),
            "height": min(table_box["height"] + 80, 4000),
        }
        await page.screenshot(path=str(img_path), clip=clip)
        print(f"  ✅ 보수 테이블 스크린샷: {img_path.name}")
        return True
    else:
        # 폴백: 전체 페이지 뷰포트 스크린샷
        await page.screenshot(path=str(img_path), full_page=False)
        print(f"  ✅ 스크린샷(폴백): {img_path.name}")
        return True


def _parse_etfcheck_holdings(text: str) -> list[dict]:
    """etfcheck /compose 페이지 텍스트에서 구성종목 5개 추출
    형식: 티커 → 통화 → 한국어명 → \\t가격 → 가격변동% → \\t비중% [→ 비중변동% (한국ETF만)]
    """
    idx = text.find("종목\t현재가\t비중")
    if idx < 0:
        return []

    chunk = text[idx + len("종목\t현재가\t비중"):]
    weight_re = re.compile(r"^\t(\d+\.\d+)%$")
    pct_re    = re.compile(r"^[+-]?\d+\.\d+%$")
    currencies = {"USD", "KRW", "EUR", "JPY", "GBP", "CAD", "TWD", "HKD", "AUD", "CHF", "SGD", "CNY"}
    stop_kws  = {"자산유형별", "섹터비중", "시가총액 규모", "국가별 비중", "회원약관", "ETF CHECK"}

    holdings = []
    last_ticker = None
    last_kr_name = None

    for line in chunk.split("\n"):
        if len(holdings) >= 5:
            break
        stripped = line.strip()
        if not stripped:
            continue
        if any(kw in stripped for kw in stop_kws):
            break

        # 비중 줄: \t숫자.숫자%
        m = weight_re.match(line)
        if m:
            weight = float(m.group(1))
            if 0.1 < weight <= 100:
                name = last_kr_name or last_ticker or "Unknown"
                holdings.append({"name": name[:30], "weight": f"{weight:.2f}%"})
            last_ticker = None
            last_kr_name = None
            continue

        # 탭으로 시작하는 가격 줄 → skip
        if line.startswith("\t"):
            continue

        # 변동률 줄 (-16.43%, 0.00% 등) → skip
        if pct_re.match(stripped):
            continue

        # 서브헤더 (전일대비) → skip
        if stripped.startswith("(") and stripped.endswith(")"):
            continue

        # 통화 코드 → skip
        if stripped in currencies:
            continue

        # 종목명/티커 처리
        has_korean = any('가' <= c <= '힣' for c in stripped)
        is_ticker = stripped.isupper() and len(stripped) <= 6 and stripped.replace(".", "").isalpha()

        if is_ticker:
            last_ticker = stripped
            last_kr_name = None
        elif has_korean:
            last_kr_name = stripped
        elif not last_ticker:
            last_ticker = stripped

    return holdings


async def scrape_etfcheck_holdings(page, compose_url: str) -> list[dict]:
    """etfcheck /compose URL에서 구성종목 추출"""
    try:
        await page.goto(compose_url, timeout=20000)
        await page.wait_for_load_state("networkidle", timeout=15000)
        await asyncio.sleep(2)
        text = await page.evaluate("() => document.body.innerText")
        return _parse_etfcheck_holdings(text)
    except Exception as e:
        print(f"    구성종목 수집 실패: {e}")
        return []


async def scrape_kr_etf_detail(page, sotCd: str, name: str | None = None) -> dict:
    """funetf 상세페이지에서 운용사/투자위험/상장일/순자산/펀드특징·운용전략/구성종목 추출"""
    result = {"운용사": None, "투자위험": None, "상장일": None, "순자산": None, "특징": None, "holdings": []}

    link = await find_kr_product_link(page, sotCd, name)
    if not link:
        return result

    await page.goto(link)
    await page.wait_for_load_state("networkidle", timeout=20000)
    await asyncio.sleep(2)
    # 스크롤해서 lazy-load 섹션 (보유비중 등) 로딩
    await page.evaluate("window.scrollTo(0, document.body.scrollHeight)")
    await asyncio.sleep(2)

    text = await page.evaluate("() => document.body.innerText")
    m = re.search(r"운용사\n(.+)", text)
    if m: result["운용사"] = m.group(1).strip()
    m = re.search(r"투자위험\n(.+)", text)
    if m: result["투자위험"] = m.group(1).strip()
    m = re.search(r"순자산\n(.+)", text)
    if m: result["순자산"] = m.group(1).strip()
    m = re.search(r"상장일\n(\d{4}\.\d{2}\.\d{2})", text)
    if m: result["상장일"] = m.group(1)

    try:
        await page.locator("#btnInfoDetail").click(timeout=5000)
        await asyncio.sleep(1.5)
        feature_text = await page.evaluate("""() => {
            const parts = [];
            for (const id of ['divDiscription1', 'divDiscription2']) {
                const el = document.getElementById(id);
                if (el) parts.push(el.innerText);
            }
            return parts.join(' ');
        }""")
        if feature_text:
            feature_text = feature_text.replace("펀드특징", "").replace("운용전략", "")
            result["특징"] = summarize(feature_text, 50)
    except Exception:
        pass

    return result


async def download_kr_pdf(page, context, sotCd: str, save_path: Path, name: str | None = None) -> bool:
    first_link = await find_kr_product_link(page, sotCd, name)

    if not first_link:
        print(f"  ❌ 상품 링크 없음")
        return False

    await page.goto(first_link)
    await page.wait_for_load_state("networkidle", timeout=20000)
    await asyncio.sleep(3)

    dl_selector = 'div.prd-11-dcmt-item:has(p.prd-11-dcmt-title:text("투자설명서")) button:has-text("다운로드")'
    try:
        btn = page.locator(dl_selector)
        if await btn.count() > 0:
            async with page.expect_download(timeout=30000) as dl_info:
                await btn.first.click()
            dl = await dl_info.value
            await dl.save_as(str(save_path))
            print(f"  ✅ PDF 저장: {save_path.name}")
            return True
    except Exception as e:
        print(f"  버튼 클릭 오류: {e}")

    try:
        fileurl = await page.evaluate('''() => {
            const items = document.querySelectorAll(".prd-11-dcmt-item");
            for (const item of items) {
                const title = item.querySelector(".prd-11-dcmt-title");
                if (title && title.textContent?.trim() === "투자설명서") {
                    const btn = item.querySelector("button.download");
                    return btn ? btn.getAttribute("data-fileurl") : null;
                }
            }
            return null;
        }''')
        if fileurl:
            full_url = f"https://www.funetf.co.kr{fileurl}"
            resp = requests.get(full_url, timeout=30)
            if resp.ok and len(resp.content) > 1000:
                save_path.write_bytes(resp.content)
                print(f"  ✅ 저장(fallback): {save_path.name}")
                return True
    except Exception as e:
        print(f"  fallback 오류: {e}")

    return False


def find_fee_page(pdf_path: Path) -> int | None:
    keywords_section = [
        "집합투자기구에 부과되는 보수",
        "(1)투자신탁", "(1) 투자신탁",
        "투자신탁 관련보수", "투자신탁 관련 보수",
    ]
    keywords_fee = [
        "증권거래비용", "증권 거래비용",
        "총보수", "총 보수",
        "합성총보수", "합성 총보수",
    ]
    candidates = []
    fallback_pg = None
    try:
        with pdfplumber.open(pdf_path) as pdf:
            total = len(pdf.pages)
            print(f"  PDF 총 {total}페이지")
            for pg_num, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                has_section = any(kw in text for kw in keywords_section)
                has_fee     = any(kw in text for kw in keywords_fee)
                if has_section and has_fee:
                    candidates.append(pg_num)
                elif has_section and pg_num + 1 < total:
                    next_text = pdf.pages[pg_num + 1].extract_text() or ""
                    if any(kw in next_text for kw in keywords_fee):
                        candidates.append(pg_num)
                if fallback_pg is None:
                    has_synth = any(k in text for k in ["합성총보수", "합성 총보수", "총보수·비용",
                                                        "총보수·비", "총보수･비용", "총보수･비"])
                    has_trade = any(k in text for k in ["증권거래비용", "증권 거래비용"])
                    if has_synth and has_trade:
                        fallback_pg = pg_num
        if candidates:
            pg = candidates[-1]
            print(f"  → 보수 섹션: {pg+1}페이지")
            return pg
        if fallback_pg is not None:
            print(f"  → 보수 섹션(fallback): {fallback_pg+1}페이지")
            return fallback_pg
    except Exception as e:
        print(f"  페이지 탐색 오류: {e}")
    return None


def capture_fee_table(pdf_path: Path, pg_num: int, out_img: Path, zoom: float = 2.5) -> Path:
    doc = fitz.open(str(pdf_path))
    section_keywords = [
        "집합투자기구에 부과되는 보수", "집합투자기구에 부과",
        "집합투자기구에부과되는보수", "집합투자기구에부과",
        "나. 집합투자기구",
    ]
    start_y_pdf = None
    page0 = doc[pg_num]
    for kw in section_keywords:
        hits = page0.search_for(kw)
        if hits:
            start_y_pdf = hits[0].y0
            break

    pages_needed = [pg_num]
    if pg_num + 1 < len(doc):
        next_text = doc[pg_num + 1].get_text()
        next_lines = next_text.split("\n")[:30]
        if any(kw in "\n".join(next_lines) for kw in ["지급시기", "증권거래비용", "합성총보수"]):
            pages_needed.append(pg_num + 1)

    page_images = []
    for pn in pages_needed:
        p = doc[pn]
        pix = p.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
        page_images.append(Image.open(BytesIO(pix.tobytes("png"))))
    doc.close()

    if len(page_images) > 1:
        total_h = sum(img.height for img in page_images)
        full_img = Image.new("RGB", (page_images[0].width, total_h), "white")
        y_off = 0
        for img in page_images:
            full_img.paste(img, (0, y_off))
            y_off += img.height
    else:
        full_img = page_images[0]

    img_w, img_h = full_img.width, full_img.height
    single_page_img_h = page_images[0].height
    single_page_img_w = page_images[0].width

    try:
        with pdfplumber.open(pdf_path) as pdf:
            pl_page0 = pdf.pages[pg_num]
            pw0 = pl_page0.width
            ph0 = pl_page0.height
            sx = single_page_img_w / pw0
            sy = single_page_img_h / ph0

            x0_pdf, x1_pdf = None, None
            end_y_img = None
            fee_kws_cap = ["총보수", "총 보수", "증권거래비용", "합성총보수", "피투자"]
            num_re_cap  = re.compile(r"\d{1,2}\.\d{2,4}")

            for i, pn in enumerate(pages_needed):
                pl_page = pdf.pages[pn]
                y_offset_img = sum(page_images[j].height for j in range(i))
                for tbl in pl_page.find_tables():
                    ext = tbl.extract()
                    if not ext: continue
                    flat = " ".join(str(c) for row in ext for c in row if c)
                    if not any(kw in flat for kw in fee_kws_cap): continue
                    total_valid = sum(
                        1 for row in ext for cell in row
                        for m in num_re_cap.findall(str(cell or ""))
                        if 0.0001 <= float(m) <= 9.9999
                    )
                    if total_valid < 2: continue
                    bbox = tbl.bbox
                    if x0_pdf is None:
                        x0_pdf = bbox[0]; x1_pdf = bbox[2]
                    end_y_img = y_offset_img + int(bbox[3] * sy) + 20

            if x0_pdf is None:
                x0_pdf, x1_pdf = 50, pw0 - 50

            cx0 = max(0, int(x0_pdf * sx) - 10)
            cx1 = min(img_w, int(x1_pdf * sx) + 10)
            cy0 = max(0, int((start_y_pdf - 20) * sy)) if start_y_pdf else 0
            cy1 = end_y_img if end_y_img else img_h
            if cy1 <= cy0:
                cy1 = img_h

            cropped = full_img.crop((cx0, cy0, cx1, cy1))
    except Exception as e:
        print(f"  크롭 오류: {e}")
        cropped = full_img

    cropped.save(str(out_img), "PNG")
    print(f"  ✅ 이미지 저장: {out_img.name}")
    return out_img


def find_fund_fee_page(pdf_path: Path) -> int | None:
    """펀드 투자설명서에서 '나. 집합투자기구에 부과되는 보수 및 비용' 본문 페이지 반환.
    경로: 제2부 → 13. 보수 및 수수료에 관한 사항 → 나. 집합투자기구에 부과되는 보수 및 비용
    목차(1~10페이지 내외)를 건너뛰고 제2부 본문에서 탐색.
    """
    # '집합투자기구에 부과되는 보수 및 비용' 섹션 키워드 (나./다./마. 등 다양한 접두사 허용)
    target_kws = [
        "나. 집합투자기구에 부과되는 보수 및 비용",
        "다. 집합투자기구에 부과되는 보수 및 비용",
        "마. 집합투자기구에 부과되는 보수 및 비용",
        "나.집합투자기구에 부과되는 보수 및 비용",
        "다.집합투자기구에 부과되는 보수 및 비용",
    ]
    # 제2부 시작 표시
    part2_kws = [
        "제2부. 집합투자기구에 관한 사항",
        "제2부.집합투자기구에 관한 사항",
        "제2부 집합투자기구에 관한 사항",
    ]
    # 폴백: (1) 이 투자신탁의 관련보수등 / 집합투자기구에 부과되는 보수
    fallback_kws = [
        "(1) 이 투자신탁의 관련보수등",
        "(1)이 투자신탁의 관련보수등",
        "집합투자기구에 부과되는 보수 및 비용",
        "집합투자기구에부과되는보수및비용",
    ]
    try:
        with pdfplumber.open(pdf_path) as pdf:
            total = len(pdf.pages)
            print(f"  PDF 총 {total}페이지")

            # 제2부 시작 페이지 탐색
            part2_start = 0
            for pg_num, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                if any(kw in text for kw in part2_kws):
                    # 페이지 첫 부분에 나타나야 본문 (목차 제외)
                    lines = text.split("\n")
                    for idx, line in enumerate(lines[:10]):
                        if any(kw in line for kw in part2_kws):
                            part2_start = pg_num
                            print(f"  → 제2부 시작: {pg_num+1}페이지")
                            break
                    if part2_start > 0:
                        break

            # 제2부 이후에서 '나. 집합투자기구에 부과되는 보수 및 비용' 탐색
            for pg_num in range(part2_start, total):
                text = pdf.pages[pg_num].extract_text() or ""
                if any(kw in text for kw in target_kws):
                    # 이 페이지에 실제 보수 데이터가 있는지 확인 (표 또는 다음 페이지 포함)
                    next_text = pdf.pages[pg_num + 1].extract_text() if pg_num + 1 < total else ""
                    if any(kw in text + next_text for kw in fallback_kws):
                        print(f"  → 집합투자기구 보수 섹션: {pg_num+1}페이지")
                        return pg_num

            # 폴백: '(1) 이 투자신탁의 관련보수등' 단독 탐색
            for pg_num in range(part2_start, total):
                text = pdf.pages[pg_num].extract_text() or ""
                if any(kw in text for kw in fallback_kws[:2]):
                    print(f"  → 보수 섹션(fallback): {pg_num+1}페이지")
                    return pg_num
    except Exception as e:
        print(f"  페이지 탐색 오류: {e}")
    return None


def extract_fund_class(fund_name: str) -> str | None:
    """펀드명에서 클래스명 추출.
    종류C-Pe / ClassC-P2e / _Ce / (Ce) / (종류C-e) / ]C-Pe / Cpe(퇴직연금) 등 처리.
    """
    if not fund_name:
        return None
    # 1. '종류XXX' 패턴
    m = re.search(r'종류([A-Za-z][A-Za-z0-9\-]{0,8})', fund_name)
    if m: return m.group(1)
    # 2. 'ClassXXX' 패턴
    m = re.search(r'Class([A-Za-z][A-Za-z0-9\-]{0,8})', fund_name)
    if m: return m.group(1)
    # 3. 괄호 안 '종류XXX'
    m = re.search(r'\(종류([A-Za-z][A-Za-z0-9\-]{0,8})\)', fund_name)
    if m: return m.group(1)
    # 4. XXX(퇴직연금) 또는 XXX(연금저축)
    m = re.search(r'([A-Za-z][A-Za-z0-9\-]{1,8})\((?:퇴직연금|연금저축)\)', fund_name)
    if m: return m.group(1)
    # 5. _XXX 패턴
    m = re.search(r'_([A-Za-z][A-Za-z0-9\-]{0,8})(?=[\[\(\t]|$)', fund_name)
    if m: return m.group(1)
    # 6. ]/[ 바로 뒤 영문 클래스 (끝 또는 연금 표시 전)
    m = re.search(r'[\]\[]([A-Za-z][A-Za-z0-9\-]{1,8})(?:\(퇴직연금\)|\(연금저축\)|$)', fund_name)
    if m: return m.group(1)
    m = re.search(r'[\]\[]([A-Za-z][A-Za-z0-9\-]{1,8})$', fund_name)
    if m: return m.group(1)
    # 7. 마지막 괄호 안 클래스 (H, UH 등 순수 대문자 약어 제외)
    matches = re.findall(r'\(([A-Za-z][A-Za-z0-9\-]{1,8})\)', fund_name)
    for ms in reversed(matches):
        if not re.match(r'^[A-Z]{1,3}$', ms):
            return ms
    # 8. 문자열 끝 영문 클래스
    m = re.search(r'([A-Za-z][A-Za-z0-9\-]{1,8})$', fund_name)
    if m: return m.group(1)
    return None


def capture_fund_fee_table(pdf_path: Path, pg_num: int, out_img: Path,
                           zoom: float = 2.5, class_name: str | None = None) -> Path:
    """펀드 투자설명서의 '집합투자기구에 부과되는 보수 및 비용' 첫 번째 표만 캡처.
    class_name 지정 시 해당 클래스 행에 빨간 테두리 표시.
    """
    doc = fitz.open(str(pdf_path))
    total_pages = len(doc)

    start_kws = [
        "나. 집합투자기구에 부과되는 보수 및 비용",
        "다. 집합투자기구에 부과되는 보수 및 비용",
        "마. 집합투자기구에 부과되는 보수 및 비용",
        "나.집합투자기구에 부과되는 보수 및 비용",
        "다.집합투자기구에 부과되는 보수 및 비용",
        "집합투자기구에 부과되는 보수 및 비용",
        "집합투자기구에부과되는보수및비용",
    ]
    end_kws = [
        "(2) 투자신탁 관련 비용",
        "(2)투자신탁 관련 비용",
        "(2) 이 투자신탁",
        "투자기간 1년",
        "클래스종류 투자기간",
        "다. 보수 및 지급내역",
        "다.보수 및 지급내역",
        "라. 보수 및 지급내역",
        "라.보수 및 지급내역",
        "14. 이익 배분",
    ]

    # 시작 페이지에서 heading y 좌표 탐색
    page0 = doc[pg_num]
    start_y_pdf = None
    for kw in start_kws:
        hits = page0.search_for(kw)
        if hits:
            start_y_pdf = hits[0].y0 - 10
            break

    # 최대 2페이지 탐색 (표는 1~2페이지)
    # end_y_pdf: 끝 키워드의 y좌표 (페이지 인덱스 포함)
    pages_needed = [pg_num]
    end_page_idx = 0   # pages_needed 내 인덱스
    end_y_pdf = None   # 끝 키워드 y (PDF 좌표)

    # 시작 페이지 자체에 끝 키워드가 있는지 확인 (heading 아래에 있을 때만)
    for kw in end_kws:
        hits = page0.search_for(kw)
        for h in hits:
            if start_y_pdf is None or h.y0 > start_y_pdf + 20:
                end_y_pdf = h.y0 - 5
                break
        if end_y_pdf is not None:
            break

    # 끝 키워드를 아직 못 찾았으면 다음 페이지 확인 (1장만)
    if end_y_pdf is None and pg_num + 1 < total_pages:
        next_pn = pg_num + 1
        next_page = doc[next_pn]
        found_end = False
        for kw in end_kws:
            hits = next_page.search_for(kw)
            if hits:
                kw_y = hits[0].y0
                page_h = next_page.rect.height
                if kw_y > page_h * 0.4:
                    # 후반부에 끝 키워드 → 이 페이지 포함, 거기서 자름
                    pages_needed.append(next_pn)
                    end_page_idx = 1
                    end_y_pdf = kw_y - 5
                else:
                    # 전반부 → 이 페이지 제외
                    pass
                found_end = True
                break
        if not found_end:
            pages_needed.append(next_pn)
            end_page_idx = 1

    print(f"  → 캡처 대상 페이지: {[p+1 for p in pages_needed]}")

    # 각 페이지 렌더링 (doc은 나중에 red box 그릴 때까지 열어둠)
    page_images = []
    for pn in pages_needed:
        pix = doc[pn].get_pixmap(matrix=fitz.Matrix(zoom, zoom))
        page_images.append(Image.open(BytesIO(pix.tobytes("png"))))

    # 수직 연결
    cumulative_h = [0]
    for img in page_images:
        cumulative_h.append(cumulative_h[-1] + img.height)
    total_h = cumulative_h[-1]

    full_img = Image.new("RGB", (page_images[0].width, total_h), "white")
    for i, img in enumerate(page_images):
        full_img.paste(img, (0, cumulative_h[i]))

    img_w = page_images[0].width

    # 위쪽 crop 좌표 (PDF y → pixel)
    cy0 = max(0, int((start_y_pdf or 0) * zoom) - 10)

    # 아래쪽 crop 좌표
    if end_y_pdf is not None:
        cy1 = cumulative_h[end_page_idx] + int(end_y_pdf * zoom)
    else:
        cy1 = total_h
    cy1 = min(cy1, total_h)

    doc.close()

    cropped = full_img.crop((0, cy0, img_w, cy1))

    # 클래스 행에 빨간 테두리 표시
    # PDF 구조: 첫 번째 열에 "수수료종류설명\n(C-pe)" 형태로 클래스코드가 괄호 안에 있음
    # pdfplumber로 왼쪽 열 텍스트를 줄별로 묶어, 괄호 포함 클래스명 검색
    if class_name:
        cls_norm = class_name.lower().replace('-', '').replace(' ', '')
        draw = ImageDraw.Draw(cropped)
        crop_h = cropped.height

        with pdfplumber.open(pdf_path) as pdf:
            for i, pn in enumerate(pages_needed):
                pl_page = pdf.pages[pn]
                page_w = pl_page.width

                # 왼쪽 열(40% 이내) 단어를 y좌표 5pt 버킷으로 묶기
                rows_by_y: dict[int, list] = {}
                for w in (pl_page.extract_words() or []):
                    if w['x0'] >= page_w * 0.4:
                        continue
                    bucket = round(w['top'] / 5) * 5
                    rows_by_y.setdefault(bucket, []).append(w)

                # 인접한 버킷들을 묶어 "논리 행" 구성 (15pt 이내 연속이면 같은 행)
                sorted_buckets = sorted(rows_by_y)
                logical_rows: list[list] = []
                for b in sorted_buckets:
                    if logical_rows and b - sorted_buckets[sorted_buckets.index(b) - 1] <= 15:
                        logical_rows[-1].extend(rows_by_y[b])
                    else:
                        logical_rows.append(list(rows_by_y[b]))

                for lr_words in logical_rows:
                    # 이 논리 행의 전체 텍스트에서 괄호 안 클래스코드 찾기
                    full_text = ''.join(w['text'] for w in lr_words)
                    # 괄호 안 문자 추출: (C-pe) → cpe
                    import re as _re
                    parens_matches = _re.findall(r'\(([^)]+)\)', full_text)
                    found = False
                    for m in parens_matches:
                        m_norm = m.lower().replace('-', '').replace(' ', '')
                        if m_norm == cls_norm:
                            found = True
                            break
                    # 괄호 없이 전체가 클래스명인 경우도 허용
                    if not found:
                        full_norm = full_text.lower().replace('-','').replace(' ','').replace('(','').replace(')','')
                        if full_norm == cls_norm:
                            found = True

                    if not found:
                        continue

                    # 행 전체 y범위 (위로 2줄분 ~25pt 확장하여 설명텍스트+숫자행 포함)
                    row_top = min(w['top'] for w in lr_words) - 25
                    row_bot = max(w['bottom'] for w in lr_words) + 5

                    y0_px = cumulative_h[i] + int(row_top * zoom) - cy0
                    y1_px = cumulative_h[i] + int(row_bot  * zoom) - cy0
                    if y1_px < 0 or y0_px > crop_h:
                        continue
                    y0_px = max(0, y0_px)
                    y1_px = min(crop_h - 1, y1_px)
                    draw.rectangle([2, y0_px, img_w - 3, y1_px], outline="red", width=3)
                    print(f"  → 클래스 행 표시: '{class_name}' y={y0_px}~{y1_px}px")

        del draw
    cropped.save(str(out_img), "PNG")
    print(f"  ✅ 펀드 보수표 저장: {out_img.name}")
    return out_img


def extract_fees(pdf_path: Path, pg_num: int) -> dict:
    result = {"synthetic_fee": None, "trading_cost": None, "breakdown": {}}
    num_find = re.compile(r"\d{1,2}\.\d{2,4}")
    num_re_hdr = re.compile(r"^[\d.,%-]+$")

    BREAKDOWN_KWS = {
        "집합": ["집합투자업자보수"],
        "지정": ["지정참가회사보수"],
        "신탁": ["신탁업자보수", "수탁회사보수"],
        "사무": ["일반사무관리회사보수", "사무관리회사보수"],
    }

    def is_valid_fee(s):
        try: return 0.0001 <= float(s) <= 9.9999
        except: return False

    pages_to_check = [pg_num, pg_num + 1]
    FEE_KWS = ["총보수", "총 보수", "증권거래비용", "증권 거래비용", "거래비용", "합성총보수"]
    SYNTH_KWS = ["피투자", "보수포함",
                 "총보수·비용", "총보수·비", "총보수·",
                 "총보수･비용", "총보수･비", "총보수･"]

    with pdfplumber.open(pdf_path) as pdf:
        for chk_pg in pages_to_check:
            if chk_pg >= len(pdf.pages): continue
            pl_page = pdf.pages[chk_pg]
            for tbl in pl_page.find_tables():
                ext = tbl.extract()
                if not ext: continue
                flat = " ".join(str(c) for row in ext for c in row if c)
                if not any(kw in flat for kw in FEE_KWS): continue

                max_cols = max(len(row) for row in ext)

                if max_cols <= 5:
                    for row in ext:
                        if not row or len(row) < 2: continue
                        key = str(row[0] or "").replace("\n", "").strip()
                        val_cells = [str(c or "") for c in row[1:] if c]
                        nums = [n for cell in val_cells for n in num_find.findall(cell) if is_valid_fee(n)]
                        if not nums: continue
                        if result["synthetic_fee"] is None and any(k in key for k in
                                ["총보수·비용","총보수·비","합성총보수","총 보수·비용","보수·비용","총보수･비용","총보수･비"]):
                            result["synthetic_fee"] = nums[0]
                        if result["trading_cost"] is None and any(k in key for k in
                                ["증권거래비용","증권 거래비용","거래비용"]):
                            result["trading_cost"] = nums[0]
                    if result["synthetic_fee"] and result["trading_cost"]:
                        return result
                else:
                    col_synth, col_trade = None, None
                    col_breakdown: dict = {}
                    col_concat: dict = {}
                    for row in ext:
                        if not row: continue
                        for ci, cell in enumerate(row):
                            t = str(cell or "").replace("\n", "").strip()
                            if t and not num_re_hdr.match(t):
                                col_concat[ci] = col_concat.get(ci, "") + t
                    for ci, t in col_concat.items():
                        if col_synth is None and any(k in t for k in SYNTH_KWS):
                            col_synth = ci
                        if col_trade is None and any(k in t for k in ["증권거래비용", "거래비용"]):
                            col_trade = ci
                        for bkey, kws in BREAKDOWN_KWS.items():
                            if bkey not in col_breakdown and any(k in t.replace(" ", "") for k in kws):
                                col_breakdown[bkey] = ci

                    flat_nums = []
                    if col_synth is not None:
                        flat_nums = [n for row in ext for ci, cell in enumerate(row)
                                     if ci == col_synth for n in num_find.findall(str(cell or ""))
                                     if is_valid_fee(n)]
                        # 합성총보수 헤더가 병합셀로 인해 값이 인접 열에 있을 수 있음
                        if not flat_nums:
                            for adj in range(1, 4):
                                for col_try in [col_synth - adj, col_synth + adj]:
                                    if col_try < 0: continue
                                    cand = [n for row in ext for ci, cell in enumerate(row)
                                            if ci == col_try for n in num_find.findall(str(cell or ""))
                                            if is_valid_fee(n)]
                                    if cand:
                                        col_synth = col_try
                                        flat_nums = cand
                                        break
                                if flat_nums:
                                    break

                    if flat_nums:
                        # 분해 항목은 합성총보수와 같은 데이터 행에서 추출
                        for row in ext:
                            if col_synth >= len(row): continue
                            row_nums = [n for n in num_find.findall(str(row[col_synth] or "")) if is_valid_fee(n)]
                            if not row_nums: continue
                            for bkey, ci in col_breakdown.items():
                                if bkey in result["breakdown"] or ci >= len(row): continue
                                bnums = [n for n in num_find.findall(str(row[ci] or "")) if is_valid_fee(n)]
                                if bnums:
                                    result["breakdown"][bkey] = bnums[0]

                        if col_synth != col_trade and col_trade is not None:
                            trade_nums = [n for row in ext for ci, cell in enumerate(row)
                                          if ci == col_trade for n in num_find.findall(str(cell or ""))
                                          if is_valid_fee(n)]
                            if trade_nums and flat_nums[-1] != trade_nums[-1]:
                                result["synthetic_fee"] = flat_nums[-1]
                                result["trading_cost"]  = trade_nums[-1]
                                return result
                        result["synthetic_fee"] = flat_nums[-1]
                    else:
                        ROW_SYNTH = ["총보수·비용","총보수·비","총보수･비용","총보수･비",
                                     "합성총보수","합성 총보수","총 보수·비용","보수·비용"]
                        ROW_TRADE = ["증권거래비용","증권 거래비용","거래비용"]
                        for row in ext:
                            if not row or len(row) < 2: continue
                            key = str(row[0] or "").replace("\n","").strip()
                            key_norm = key.replace(" ", "")
                            val_cells = [str(c or "") for c in row[1:] if c]
                            nums = [n for cell in val_cells for n in num_find.findall(cell) if is_valid_fee(n)]
                            if not nums: continue
                            if result["synthetic_fee"] is None and any(k in key for k in ROW_SYNTH):
                                result["synthetic_fee"] = nums[0]
                            if result["trading_cost"] is None and any(k in key for k in ROW_TRADE):
                                result["trading_cost"] = nums[0]
                            for bkey, kws in BREAKDOWN_KWS.items():
                                if bkey not in result["breakdown"] and any(k == key_norm for k in kws):
                                    result["breakdown"][bkey] = nums[0]

        # Text fallback
        for chk_pg in pages_to_check:
            if chk_pg >= len(pdf.pages): continue
            text = pdf.pages[chk_pg].extract_text() or ""
            for line in text.split("\n"):
                line = line.strip()
                nums = [n for n in num_find.findall(line) if is_valid_fee(n)]
                if not nums: continue
                if result["synthetic_fee"] is None and any(k in line for k in
                        ["합성총보수","합성 총보수","총보수·비용","총보수·비","총보수･비용","총보수･비","피투자","보수포함"]):
                    result["synthetic_fee"] = nums[-1]
                if result["trading_cost"] is None and any(k in line for k in ["증권거래비용","증권 거래비용"]):
                    result["trading_cost"] = nums[-1]

    return result


# ═══════════════════════════════════════════════════════════════════════════
# ── 미국 ETF 스크래핑 ────────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════

async def scrape_us_etf(page, item: dict) -> dict:
    ticker = item["ticker"]
    etf_id = item.get("etfcheck_id")
    result = {"ticker": ticker, "name": item["name"],
              "설정일": None, "시가총액": None, "TER": None, "운용사": None, "특징": None,
              "holdings": [], "img": None}

    if not etf_id:
        print(f"  ⚠️ {ticker}: etfcheck ID 없음")
        return result

    url = f"https://www.etfcheck.co.kr/mobile/global/etpitem/{etf_id}/basic/개요"
    await page.goto(url, timeout=25000)
    await page.wait_for_load_state("networkidle", timeout=20000)
    await asyncio.sleep(3)

    text = await page.evaluate("() => document.body.innerText")

    m = re.search(r"설정일\s+(\d{4}\.\d{2}\.\d{2})", text)
    if m: result["설정일"] = m.group(1)
    m = re.search(r"시가총액\s+(\$\s*[\d,]+\.?\d*[A-Z]?)", text)
    if m: result["시가총액"] = m.group(1).strip()
    m = re.search(r"TER\s+([\d.]+%)", text)
    if m: result["TER"] = m.group(1)
    m = re.search(r"운용사\s+([^\n]+)", text)
    if m: result["운용사"] = m.group(1).strip()

    # 투자전략(한글 번역) 요약
    idx = text.find("투자전략")
    if idx >= 0:
        chunk = text[idx + len("투자전략"):]
        end_idx = chunk.find("※ 한글 설명")
        if end_idx >= 0:
            chunk = chunk[:end_idx]
        paragraphs = [p.strip() for p in chunk.split("\n\n") if p.strip()]
        if paragraphs:
            result["특징"] = summarize(paragraphs[-1], 50)

    print(f"  ✅ {ticker}: 설정일={result['설정일']}, 시가총액={result['시가총액']}, TER={result['TER']}, 운용사={result['운용사']}")

    # 스크린샷 (기본정보 섹션으로 스크롤)
    img_path = IMG_DIR / f"{ticker}_etfcheck.png"
    try:
        await page.evaluate("""() => {
            const els = document.querySelectorAll('*');
            for (const el of els) {
                if (el.children.length === 0 &&
                    (el.textContent.trim() === 'ETF 개요' ||
                     el.textContent.trim() === '설정일')) {
                    el.scrollIntoView({behavior: 'instant', block: 'start'});
                    return true;
                }
            }
            window.scrollBy(0, 480);
            return false;
        }""")
        await asyncio.sleep(1)
    except: pass
    await page.screenshot(path=str(img_path), full_page=False)
    result["img"] = img_path
    print(f"  ✅ 스크린샷: {img_path.name}")

    # 구성종목: etfcheck /compose
    if etf_id:
        compose_url = f"https://www.etfcheck.co.kr/mobile/global/etpitem/{etf_id}/compose"
        result["holdings"] = await scrape_etfcheck_holdings(page, compose_url)
        if result["holdings"]:
            print(f"  ▶ 구성종목 {len(result['holdings'])}개 수집 (etfcheck)")

    return result


# ═══════════════════════════════════════════════════════════════════════════
# ── PPT 생성 (space_etf_report.py 와 동일) ──────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════

TITLE_COLOR  = RGBColor(0x1F, 0x49, 0x7D)
HEADER_COLOR = RGBColor(0x1F, 0x49, 0x7D)
ROW_ALT      = RGBColor(0xE9, 0xF0, 0xF8)
ROW_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
FONT_KR      = "맑은 고딕"


def set_cell_text(cell, text, font_size=11, bold=False,
                  bg_color=None, align=PP_ALIGN.CENTER, color=None):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = FONT_KR
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    if bg_color:
        from lxml import etree
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        solidFill = etree.SubElement(tcPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
        srgbClr = etree.SubElement(solidFill, "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
        srgbClr.set("val", str(bg_color))


def format_kr_total_fee(r) -> str:
    sf = r.get("synthetic_fee")
    if not sf:
        return "N/A"
    bd = r.get("breakdown") or {}
    listed = r.get("상장일")
    is_recent = False
    if listed:
        try:
            listed_dt = datetime.strptime(listed, "%Y.%m.%d")
            is_recent = (datetime.now() - listed_dt).days < 365
        except ValueError:
            pass
    if is_recent and all(k in bd for k in ("집합", "지정", "신탁", "사무")):
        return (f"{sf}%(집합 {bd['집합']}%, 지정 {bd['지정']}%, "
                f"신탁 {bd['신탁']}%, 사무 {bd['사무']}%)")
    return f"{sf}%"


def add_summary_table(prs, kr_results, us_results, fund_results=None):
    fund_results = fund_results or []
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width
    txb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), sw - Inches(0.8), Inches(0.6))
    run = txb.text_frame.paragraphs[0].add_run()
    run.text = "ETF / 펀드 요약"
    run.font.name = FONT_KR; run.font.size = Pt(22); run.font.bold = True
    run.font.color.rgb = TITLE_COLOR

    headers = ["상품명", "코드", "설정/상장일", "순자산", "투자위험도", "자산운용사", "총보수", "증권거래비용", "특징"]
    col_w = [Inches(1.9), Inches(0.7), Inches(0.9), Inches(1.1), Inches(0.9),
             Inches(1.1), Inches(2.6), Inches(1.0), Inches(2.73)]
    n_rows = 1 + len(kr_results) + len(us_results) + len(fund_results)
    tbl = slide.shapes.add_table(n_rows, len(headers), Inches(0.3), Inches(1.0),
                                  sum(col_w), Inches(0.38 * n_rows)).table
    for ci, w in enumerate(col_w): tbl.columns[ci].width = w
    for ci, h in enumerate(headers):
        set_cell_text(tbl.cell(0, ci), h, font_size=11, bold=True,
                      bg_color=HEADER_COLOR, color=RGBColor(0xFF,0xFF,0xFF))

    ri = 1
    for r in kr_results:
        bg = ROW_ALT if ri % 2 == 0 else ROW_WHITE
        tc = f"{r['trading_cost']}%" if r.get("trading_cost") else "N/A"
        set_cell_text(tbl.cell(ri, 0), r["name"],            font_size=10, bg_color=bg, align=PP_ALIGN.LEFT)
        set_cell_text(tbl.cell(ri, 1), r["sotCd"],           font_size=10, bg_color=bg)
        set_cell_text(tbl.cell(ri, 2), r.get("상장일") or "N/A",  font_size=10, bg_color=bg)
        set_cell_text(tbl.cell(ri, 3), r.get("순자산") or "N/A",  font_size=10, bg_color=bg)
        set_cell_text(tbl.cell(ri, 4), r.get("투자위험") or "N/A", font_size=10, bg_color=bg)
        set_cell_text(tbl.cell(ri, 5), r.get("운용사") or "N/A",  font_size=10, bg_color=bg)
        set_cell_text(tbl.cell(ri, 6), format_kr_total_fee(r),  font_size=9,  bg_color=bg)
        set_cell_text(tbl.cell(ri, 7), tc,                       font_size=10, bg_color=bg)
        set_cell_text(tbl.cell(ri, 8), r.get("특징") or "N/A",   font_size=9,  bg_color=bg, align=PP_ALIGN.LEFT)
        ri += 1

    for r in us_results:
        bg = ROW_ALT if ri % 2 == 0 else ROW_WHITE
        set_cell_text(tbl.cell(ri, 0), r["name"],            font_size=10, bg_color=bg, align=PP_ALIGN.LEFT)
        set_cell_text(tbl.cell(ri, 1), r["ticker"],          font_size=10, bg_color=bg)
        set_cell_text(tbl.cell(ri, 2), r.get("설정일") or "N/A",  font_size=10, bg_color=bg)
        set_cell_text(tbl.cell(ri, 3), r.get("시가총액") or "N/A", font_size=10, bg_color=bg)
        set_cell_text(tbl.cell(ri, 4), "-",                   font_size=10, bg_color=bg)
        set_cell_text(tbl.cell(ri, 5), r.get("운용사") or "N/A", font_size=10, bg_color=bg)
        set_cell_text(tbl.cell(ri, 6), r.get("TER") or "N/A",   font_size=10, bg_color=bg)
        set_cell_text(tbl.cell(ri, 7), "-",                   font_size=10, bg_color=bg)
        set_cell_text(tbl.cell(ri, 8), r.get("특징") or "N/A", font_size=9,  bg_color=bg, align=PP_ALIGN.LEFT)
        ri += 1

    for r in fund_results:
        bg = ROW_ALT if ri % 2 == 0 else ROW_WHITE
        tc = f"{r['trading_cost']}%" if r.get("trading_cost") else "-"
        set_cell_text(tbl.cell(ri, 0), r["name"],            font_size=10, bg_color=bg, align=PP_ALIGN.LEFT)
        set_cell_text(tbl.cell(ri, 1), "-",                  font_size=10, bg_color=bg)
        set_cell_text(tbl.cell(ri, 2), r.get("상장일") or "N/A",  font_size=10, bg_color=bg)
        set_cell_text(tbl.cell(ri, 3), r.get("순자산") or "N/A",  font_size=10, bg_color=bg)
        set_cell_text(tbl.cell(ri, 4), r.get("투자위험") or "N/A", font_size=10, bg_color=bg)
        set_cell_text(tbl.cell(ri, 5), r.get("운용사") or "N/A",  font_size=10, bg_color=bg)
        set_cell_text(tbl.cell(ri, 6), format_kr_total_fee(r),  font_size=9,  bg_color=bg)
        set_cell_text(tbl.cell(ri, 7), tc,                       font_size=10, bg_color=bg)
        set_cell_text(tbl.cell(ri, 8), r.get("특징") or "N/A",   font_size=9,  bg_color=bg, align=PP_ALIGN.LEFT)
        ri += 1


def add_holdings_panel(slide, sw, sh, holdings: list, top):
    """슬라이드 오른쪽에 상위 구성종목 미니 테이블 추가"""
    if not holdings:
        return
    panel_w = Inches(3.7)
    panel_x = sw - panel_w - Inches(0.15)

    txb = slide.shapes.add_textbox(panel_x, top, panel_w, Inches(0.32))
    r = txb.text_frame.paragraphs[0].add_run()
    r.text = "상위 구성종목"
    r.font.name = FONT_KR; r.font.size = Pt(10); r.font.bold = True
    r.font.color.rgb = TITLE_COLOR

    tbl_top = top + Inches(0.33)
    n = min(len(holdings), 5)
    tbl = slide.shapes.add_table(n + 1, 3, panel_x, tbl_top, panel_w, Inches(0.32 * (n + 1))).table
    tbl.columns[0].width = Inches(0.32)
    tbl.columns[1].width = Inches(2.58)
    tbl.columns[2].width = Inches(0.8)

    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    set_cell_text(tbl.cell(0, 0), "#",    font_size=9, bold=True, bg_color=HEADER_COLOR, color=WHITE)
    set_cell_text(tbl.cell(0, 1), "종목명", font_size=9, bold=True, bg_color=HEADER_COLOR, color=WHITE, align=PP_ALIGN.LEFT)
    set_cell_text(tbl.cell(0, 2), "비중",  font_size=9, bold=True, bg_color=HEADER_COLOR, color=WHITE)

    for i, h in enumerate(holdings[:5]):
        bg = ROW_ALT if i % 2 == 0 else ROW_WHITE
        set_cell_text(tbl.cell(i + 1, 0), str(i + 1),        font_size=9, bg_color=bg)
        set_cell_text(tbl.cell(i + 1, 1), h.get("name", ""), font_size=9, bg_color=bg, align=PP_ALIGN.LEFT)
        set_cell_text(tbl.cell(i + 1, 2), h.get("weight", ""), font_size=9, bg_color=bg)


def add_kr_etf_slide(prs, item, fees, img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width; sh = prs.slide_height
    txb = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), sw-Inches(0.6), Inches(0.45))
    run = txb.text_frame.paragraphs[0].add_run()
    run.text = item["name"]
    run.font.name = FONT_KR; run.font.size = Pt(18); run.font.bold = True
    run.font.color.rgb = TITLE_COLOR

    sf = f"{fees.get('synthetic_fee')}%" if fees.get("synthetic_fee") else "N/A"
    tc = f"{fees.get('trading_cost')}%"  if fees.get("trading_cost")  else "N/A"
    txb2 = slide.shapes.add_textbox(Inches(0.3), Inches(0.62), sw-Inches(0.6), Inches(0.35))
    run2 = txb2.text_frame.paragraphs[0].add_run()
    run2.text = f"총보수·비용(피투자포함): {sf}    증권거래비용: {tc}"
    run2.font.name = FONT_KR; run2.font.size = Pt(11)
    run2.font.color.rgb = RGBColor(0x44,0x44,0x44)

    holdings = item.get("holdings", [])
    has_holdings = len(holdings) > 0
    img_top = Inches(1.05)

    if img_path and Path(img_path).exists():
        img = Image.open(img_path)
        iw, ih = img.size
        avail_w = (sw - Inches(4.1)) if has_holdings else (sw - Inches(0.6))
        avail_h = sh - Inches(1.15)
        scale = min(avail_w / (iw * 9144), avail_h / (ih * 9144))
        disp_w = int(iw * 9144 * scale); disp_h = int(ih * 9144 * scale)
        img_x = Inches(0.3) if has_holdings else (sw - disp_w) // 2
        slide.shapes.add_picture(str(img_path), img_x, img_top, disp_w, disp_h)

    if has_holdings:
        add_holdings_panel(slide, sw, sh, holdings, img_top)


def add_us_etf_slide(prs, us_result):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width; sh = prs.slide_height
    txb = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), sw-Inches(0.6), Inches(0.45))
    run = txb.text_frame.paragraphs[0].add_run()
    run.text = f"{us_result['ticker']}  ({us_result['name']})"
    run.font.name = FONT_KR; run.font.size = Pt(16); run.font.bold = True
    run.font.color.rgb = TITLE_COLOR

    txb2 = slide.shapes.add_textbox(Inches(0.3), Inches(0.62), sw-Inches(0.6), Inches(0.35))
    run2 = txb2.text_frame.paragraphs[0].add_run()
    run2.text = f"설정일: {us_result['설정일'] or 'N/A'}    시가총액: {us_result['시가총액'] or 'N/A'}    TER: {us_result['TER'] or 'N/A'}"
    run2.font.name = FONT_KR; run2.font.size = Pt(11)
    run2.font.color.rgb = RGBColor(0x44,0x44,0x44)

    holdings = us_result.get("holdings", [])
    has_holdings = len(holdings) > 0
    img_top = Inches(1.05)

    img_path = us_result.get("img")
    if img_path and Path(img_path).exists():
        img = Image.open(img_path)
        iw, ih = img.size
        avail_w = (sw - Inches(4.1)) if has_holdings else (sw - Inches(0.6))
        avail_h = sh - Inches(1.15)
        scale = min(avail_w / (iw * 9144), avail_h / (ih * 9144))
        disp_w = int(iw * 9144 * scale); disp_h = int(ih * 9144 * scale)
        img_x = Inches(0.3) if has_holdings else (sw - disp_w) // 2
        slide.shapes.add_picture(str(img_path), img_x, img_top, disp_w, disp_h)
    else:
        txb3 = slide.shapes.add_textbox(Inches(0.3), Inches(1.3), sw-Inches(0.6), Inches(2))
        txb3.text_frame.paragraphs[0].add_run().text = "etfcheck.co.kr 데이터 없음 (N/A)"

    if has_holdings:
        add_holdings_panel(slide, sw, sh, holdings, img_top)


def build_ppt(kr_items_with_fees, us_results, fund_results=None, title_suffix=""):
    fund_results = fund_results or []
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    add_summary_table(prs, kr_items_with_fees, us_results, fund_results)
    for item_fees in kr_items_with_fees:
        add_kr_etf_slide(prs, item_fees, item_fees, item_fees.get("img"))
    for r in us_results:
        add_us_etf_slide(prs, r)
    for r in fund_results:
        add_kr_etf_slide(prs, r, r, r.get("img"))
    out = OUT_DIR / f"etf_report_{TIMESTAMP}.pptx"
    prs.save(str(out))
    print(f"\n✅ PPT 저장: {out}")
    return out


# ═══════════════════════════════════════════════════════════════════════════
# ── MAIN ────────────────────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════

async def main():
    print("=" * 60)
    print("ETF 보고서 생성기")
    print("=" * 60)

    kr_items, us_items_raw, fund_items = load_config("config.txt")
    kr_items_with_fees = []
    us_results = []
    fund_results = []

    async with async_playwright() as p:
        # ── 브라우저 1: funetf (한국 ETF + 일반 펀드) ───────────────────
        browser_kr = await p.chromium.launch(headless=True)
        context_kr = await browser_kr.new_context(
            user_agent='Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36',
            accept_downloads=True, viewport={"width": 1400, "height": 900},
        )
        page_kr = await context_kr.new_page()
        print("\n[funetf 로그인]")
        await login(page_kr)

        for item in kr_items:
            name  = item["name"]
            sotCd = item["sotCd"]
            safe  = name.replace(" ", "_").replace("/", "_")
            pdf_path = PDF_DIR / f"{safe}.pdf"
            img_path = IMG_DIR / f"{safe}_fee.png"

            print(f"\n{'='*50}\n▶ {name} ({sotCd})\n{'='*50}")

            detail = await scrape_kr_etf_detail(page_kr, sotCd, name if name != sotCd else None)
            print(f"  ▶ 운용사={detail['운용사']}  상장일={detail['상장일']}  "
                  f"순자산={detail['순자산']}  투자위험={detail['투자위험']}")

            if pdf_path.exists() and pdf_path.stat().st_size > 10000:
                print(f"  PDF 재사용: {pdf_path.name}")
            else:
                ok = await download_kr_pdf(page_kr, context_kr, sotCd, pdf_path, name if name != sotCd else None)
                if not ok:
                    kr_items_with_fees.append({**item, **detail, "synthetic_fee": None, "trading_cost": None,
                                                "breakdown": {}, "img": None})
                    continue

            pg_num = find_fee_page(pdf_path)
            if pg_num is None:
                kr_items_with_fees.append({**item, **detail, "synthetic_fee": None, "trading_cost": None,
                                            "breakdown": {}, "img": None})
                continue

            capture_fee_table(pdf_path, pg_num, img_path)
            fees = extract_fees(pdf_path, pg_num)
            print(f"  ▶ 합성총보수: {fees['synthetic_fee']}%  증권거래비용: {fees['trading_cost']}%  "
                  f"분해: {fees['breakdown']}")

            kr_items_with_fees.append({
                **item,
                **detail,
                "synthetic_fee": fees["synthetic_fee"],
                "trading_cost":  fees["trading_cost"],
                "breakdown":     fees["breakdown"],
                "img": str(img_path) if img_path.exists() else None,
            })

        # ── 일반 펀드 처리 (funetf, 같은 브라우저) ──────────────────────
        if fund_items:
            print(f"\n\n{'='*50}\n▶ 일반 펀드\n{'='*50}")
            for item in fund_items:
                name = item["name"]
                fund_id = item.get("_fund_id")
                safe = re.sub(r'[\\/:*?"<>|]', '_', fund_id or name)
                pdf_path = PDF_DIR / f"fund_{safe}.pdf"
                img_path = IMG_DIR / f"fund_{safe}_fee.png"

                # 클래스명: config의 _class 또는 config 펀드명 끝 토큰에서 추출
                # (funetf 펀드명으로 덮어쓰기 전에 먼저 추출)
                fund_class = item.get("_class") or extract_fund_class(item["name"])

                print(f"\n[펀드] {name}")

                if not fund_id:
                    print(f"  ❌ fund_id 없음 — config에 'KR코드' 입력 필요")
                    fund_results.append({**item, "synthetic_fee": None,
                                         "trading_cost": None, "breakdown": {}, "img": None})
                    continue

                detail = await scrape_fund_detail(page_kr, name, fund_id)
                if detail.get("펀드명"):
                    name = detail["펀드명"]
                print(f"  ▶ 운용사={detail['운용사']}  설정일={detail['상장일']}  "
                      f"순자산={detail['순자산']}  투자위험={detail['투자위험']}")

                # PDF 다운로드
                if pdf_path.exists() and pdf_path.stat().st_size > 10000:
                    print(f"  PDF 재사용: {pdf_path.name}")
                else:
                    ok = await download_fund_pdf(page_kr, context_kr, name, pdf_path, fund_id)
                    if not ok:
                        fund_results.append({**item, **detail, "name": name,
                                             "synthetic_fee": None, "trading_cost": None,
                                             "breakdown": {}, "img": None})
                        continue

                # 종류별 보수 섹션 탐색 및 캡처
                pg_num = find_fund_fee_page(pdf_path)
                if pg_num is None:
                    print(f"  ⚠️ 보수 섹션 못 찾음")
                    fund_results.append({**item, **detail, "name": name,
                                         "synthetic_fee": None, "trading_cost": None,
                                         "breakdown": {}, "img": None})
                    continue

                # PDF 보수 페이지에서 실제 클래스 목록 출력 (참고용)
                with pdfplumber.open(pdf_path) as _pdf:
                    all_codes = []
                    for _pn in [pg_num, pg_num + 1]:
                        if _pn < len(_pdf.pages):
                            _text = _pdf.pages[_pn].extract_text() or ""
                            all_codes += re.findall(r'\(([A-Za-z][A-Za-z0-9\-]{0,6})\)', _text)
                    all_codes = list(dict.fromkeys(all_codes))
                    print(f"  → PDF 클래스 목록: {all_codes}")

                if fund_class:
                    print(f"  → 강조 클래스: {fund_class}")
                else:
                    print(f"  ⚠️  클래스 미지정 — config에 'FUND_ID 클래스명' 형태로 입력하세요")
                capture_fund_fee_table(pdf_path, pg_num, img_path, class_name=fund_class)

                fund_results.append({
                    **item,
                    **detail,
                    "name": name,
                    "synthetic_fee": None,
                    "trading_cost":  None,
                    "breakdown":     {},
                    "img": str(img_path) if img_path.exists() else None,
                })

        await browser_kr.close()

        # ── 브라우저 2: etfcheck (한국 ETF 구성종목 + 미국 ETF 전체) ──
        browser_us = await p.chromium.launch(headless=True)
        context_us = await browser_us.new_context(
            user_agent='Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36',
            viewport={"width": 1200, "height": 900},
        )
        page_us = await context_us.new_page()

        print("\n[etfcheck 로그인]")
        await etfcheck_login(page_us)

        # 한국 ETF 구성종목 수집 (etfcheck)
        if kr_items_with_fees:
            print("\n[한국 ETF - 구성종목 (etfcheck)]")
            for item in kr_items_with_fees:
                sotCd = item["sotCd"]
                compose_url = f"https://www.etfcheck.co.kr/mobile/etpitem/{sotCd}/compose"
                holdings = await scrape_etfcheck_holdings(page_us, compose_url)
                item["holdings"] = holdings
                print(f"  {item['name']}: {len(holdings)}개")

        # 미국 ETF
        if us_items_raw:
            print(f"\n\n{'='*50}\n▶ 미국 ETF (etfcheck.co.kr)\n{'='*50}")
            print("\n  [etfcheck ID 확인 중...]")
            for item in us_items_raw:
                ticker = item["ticker"]
                preset = item.get("_preset_id")
                print(f"\n  [{ticker}]")

                if preset == "N/A":
                    item["etfcheck_id"] = None
                    print(f"    → N/A (config에 명시)")
                elif preset and preset.startswith("F"):
                    item["etfcheck_id"] = preset
                    print(f"    → ID 직접 지정: {preset}")
                else:
                    found = await find_etfcheck_id(page_us, ticker)
                    if found:
                        item["etfcheck_id"], item["name"] = found
                    else:
                        item["etfcheck_id"] = None

                result = await scrape_us_etf(page_us, item)
                us_results.append(result)

        await browser_us.close()

    # ── PPT 생성 ──────────────────────────────────────────────────────
    print("\n\n[PPT 생성]")
    ppt_path = build_ppt(kr_items_with_fees, us_results, fund_results)

    print("\n" + "=" * 60)
    print("한국 ETF 요약")
    print("-" * 60)
    for r in kr_items_with_fees:
        print(f"  {r['name']:25s} 총보수={format_kr_total_fee(r)}")

    print("\n미국 ETF 요약")
    print("-" * 60)
    for r in us_results:
        print(f"  {r['ticker']:6s} 설정일={r['설정일'] or 'N/A':12s} 시가총액={r['시가총액'] or 'N/A':15s} TER={r['TER'] or 'N/A'}")

    if fund_results:
        print("\n펀드 요약")
        print("-" * 60)
        for r in fund_results:
            print(f"  {r['name']:30s} 총보수={format_kr_total_fee(r)}")

    print(f"\n📄 결과 파일: {ppt_path}")
    try:
        input("\nEnter 키를 누르면 창이 닫힙니다...")
    except EOFError:
        pass


asyncio.run(main())
