"""
FunETF 완전 자동화 스크래퍼
- 슬라이드 17의 모든 ETF/펀드 데이터 수집
- 상장일/설정일, 자산운용사, 투자위험도, 합성총보수, 증권거래비용
- 투자설명서 PDF에서 보수 표 캡쳐 (전체 클래스)
- 결과 PPT 파일 생성
"""
import asyncio
import json
import os
import re
import sys
import requests
import pdfplumber
import pandas as pd
from io import BytesIO
from pathlib import Path
from datetime import datetime
import fitz  # PyMuPDF for PDF screenshot
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

EMAIL = "sungsoowon45@gmail.com"
PASSWORD = "!sungsoo0405"
OUTPUT_DIR = Path(r"C:\Users\jswon\Desktop\업무\funetf_scraper\output")
PDF_DIR = OUTPUT_DIR / "pdfs"
IMG_DIR = OUTPUT_DIR / "images"
OUTPUT_DIR.mkdir(exist_ok=True)
PDF_DIR.mkdir(exist_ok=True)
IMG_DIR.mkdir(exist_ok=True)

# ── 슬라이드 17 항목 목록 ──
ITEMS = [
    # (display_name, type, search_key)
    # type: 'etf' or 'fund'
    # search_key: for ETF=sotCd(6자리), for fund=shortCode or search name

    # Focus 연금 ETF
    {"name": "KODEX 미국S&P500",                         "type": "etf",  "sotCd": "379800", "search": "KODEX 미국S&P500"},
    {"name": "TIME 글로벌AI인공지능액티브",                   "type": "etf",  "sotCd": "456600", "search": "TIME 글로벌AI인공지능액티브"},
    {"name": "KODEX AI전력핵심설비",                        "type": "etf",  "sotCd": "487240", "search": "KODEX AI전력핵심설비"},
    {"name": "ACE 글로벌반도체TOP4 Plus",                   "type": "etf",  "sotCd": "446770", "search": "ACE 글로벌반도체TOP4 Plus"},
    {"name": "RISE 대형고배당10TR",                         "type": "etf",  "sotCd": "315960", "search": "RISE 대형고배당10TR"},
    {"name": "KODEX 머니마켓액티브",                         "type": "etf",  "sotCd": "488770", "search": "KODEX 머니마켓액티브"},
    {"name": "RISE 미국단기투자등급회사채액티브",                 "type": "etf",  "sotCd": "437350", "search": "RISE 미국단기투자등급회사채액티브"},
    {"name": "ACE KRX금현물",                              "type": "etf",  "sotCd": "411060", "search": "ACE KRX금현물"},
    {"name": "KODEX 미국부동산리츠(H)",                       "type": "etf",  "sotCd": "352560", "search": "KODEX 미국부동산리츠"},
    {"name": "TIGER 글로벌멀티에셋TIF액티브",                   "type": "etf",  "sotCd": "440340", "search": "TIGER 글로벌멀티에셋TIF"},
    {"name": "KODEX TRF3070",                             "type": "etf",  "sotCd": "329650", "search": "KODEX TRF3070"},

    # Focus 퇴직연금 펀드
    {"name": "삼성미국S&P500인덱스증권자투자신탁UH[주식]Cp(퇴직연금)",          "type": "fund", "code": "PAS025", "search": "삼성미국S&P500인덱스증권자투자신탁UH"},
    {"name": "미래에셋미국나스닥100인덱스증권자투자신탁(주식)(UH)C-P2",           "type": "fund", "code": "PAM078", "search": "미래에셋미국나스닥100인덱스증권자투자신탁"},
    {"name": "KCGI코리아퇴직연금증권자투자신탁[주식]C",                        "type": "fund", "code": "PAMJ03", "search": "KCGI코리아퇴직연금증권자투자신탁"},
    {"name": "NH-Amundi 하나로단기채증권투자신탁[채권]C-P2(퇴직연금)",            "type": "fund", "code": "PDNC74", "search": "NH-Amundi 하나로단기채증권투자신탁"},
    {"name": "신한MAN글로벌투자등급채권증권투자신탁(H)[채권-재간접형]C-r",           "type": "fund", "code": "PGJA46", "search": "신한MAN글로벌투자등급채권증권투자신탁"},
    {"name": "신한골드증권투자신탁제1호[주식]C-r",                            "type": "fund", "code": "PAJA24", "search": "신한골드증권투자신탁"},
    {"name": "하나글로벌리츠부동산자투자신탁[재간접형]C-P2",                      "type": "fund", "code": "PGD862", "search": "하나글로벌리츠부동산자투자신탁"},
    {"name": "NH-Amundi하나로적격TDF2040증권투자[주식혼합-재간접형]C-P2",          "type": "fund", "code": "PGNC86", "search": "NH-Amundi하나로적격TDF2040"},
    {"name": "KB온국민적격TDF2040증권자투자[주식혼합-재간접형](H)C-퇴직",           "type": "fund", "code": "PGZA34", "search": "KB온국민적격TDF2040"},
    {"name": "미래에셋전략배분적격TDF2040혼합자산자투자신탁C-P2",                  "type": "fund", "code": "PBM573", "search": "미래에셋전략배분적격TDF2040"},

    # Focus 연금저축 펀드
    {"name": "삼성미국S&P500인덱스증권자투자신탁UH[주식]C-Pe",                   "type": "fund", "code": "NAQP03", "search": "삼성미국S&P500인덱스증권자투자신탁UH"},
    {"name": "미래에셋미국나스닥100인덱스증권자투자신탁(주식)(UH)C-Pe",             "type": "fund", "code": "NAM079", "search": "미래에셋미국나스닥100인덱스증권자투자신탁"},
    {"name": "KCGI코리아증권투자신탁1호[주식]C-Pe",                           "type": "fund", "code": "NAMP94", "search": "KCGI코리아증권투자신탁"},
    {"name": "NH-Amundi 하나로단기채증권투자신탁[채권]C-P1e(연금저축)",             "type": "fund", "code": "NDNP88", "search": "NH-Amundi 하나로단기채증권투자신탁"},
    {"name": "신한MAN글로벌투자등급채권증권투자신탁(H)[채권-재간접형]C-pe",            "type": "fund", "code": "NGJA45", "search": "신한MAN글로벌투자등급채권증권투자신탁"},
    {"name": "신한골드증권투자신탁제1호[주식]C-pe",                             "type": "fund", "code": "NASJ83", "search": "신한골드증권투자신탁"},
    {"name": "하나글로벌리츠부동산자투자신탁[재간접형]C-PE",                       "type": "fund", "code": "NGDP20", "search": "하나글로벌리츠부동산자투자신탁"},
    {"name": "NH-Amundi하나로적격TDF2040증권투자[주식혼합-재간접형]C-P1e",          "type": "fund", "code": "NGUP02", "search": "NH-Amundi하나로적격TDF2040"},
    {"name": "KB온국민적격TDF2040증권자투자[주식혼합-재간접형](H)C-Pe",             "type": "fund", "code": "NGZA18", "search": "KB온국민적격TDF2040"},
    {"name": "미래에셋전략배분적격TDF2040혼합자산자투자신탁C-Pe",                   "type": "fund", "code": "NBMP57", "search": "미래에셋전략배분적격TDF2040"},
]


def fmt_date(ymd_str):
    """YYYYMMDD → YYYY.MM.DD"""
    if not ymd_str or len(str(ymd_str)) < 8:
        return str(ymd_str) if ymd_str else "-"
    s = str(ymd_str)
    return f"{s[:4]}.{s[4:6]}.{s[6:8]}"


def risk_grade_label(grade):
    mapping = {"1": "1등급(매우높음)", "2": "2등급(높음)", "3": "3등급(중간)", "4": "4등급(낮음)", "5": "5등급(매우낮음)"}
    return mapping.get(str(grade), str(grade) if grade else "-")


async def login(page):
    await page.goto("https://www.funetf.co.kr/auth")
    await page.wait_for_load_state("domcontentloaded")
    await page.fill('input[name="username"]', EMAIL)
    await page.fill('input[name="password"]', PASSWORD)
    await page.locator('button:has-text("로그인")').last.click()
    await page.wait_for_url("https://www.funetf.co.kr/", timeout=15000)
    print("✅ 로그인 성공")


async def search_fund_etf(page, search_term):
    """검색 API로 펀드/ETF 정보 조회"""
    import urllib.parse
    encoded = urllib.parse.quote(search_term)
    url = f"https://www.funetf.co.kr/api/public/main/search/all?schVal={encoded}&reSchVal=&schKeyword="
    await page.goto(url)
    await asyncio.sleep(0.5)
    try:
        data = json.loads(await page.inner_text("body"))
        etf_list = data.get("etfList", {}).get("content", [])
        fund_list = data.get("fundList", {}).get("content", [])
        return etf_list, fund_list
    except Exception as e:
        print(f"  검색 오류: {e}")
        return [], []


async def get_detail_page_data(page, url):
    """상세 페이지 로드 후 API 호출 캡쳐 + 파일 URL 수집"""
    api_calls = {}
    file_urls = {}  # R1=투자설명서, R2=간이투자설명서, etc.

    async def on_response(response):
        resp_url = response.url
        ct = response.headers.get("content-type", "")
        if "funetf.co.kr/api" in resp_url and "json" in ct:
            endpoint = resp_url.split("/api/public/product/view/")[1].split("?")[0] if "/api/public/product/view/" in resp_url else ""
            if endpoint:
                try:
                    body = await response.json()
                    api_calls[endpoint] = body
                except:
                    pass

    page.on("response", on_response)
    try:
        await page.goto(url)
        await page.wait_for_load_state("networkidle", timeout=20000)
        await asyncio.sleep(3)

        # 페이지 DOM에서 파일 URL 추출 (R1=투자설명서, R2=간이, R3=핵심정보)
        dom_file_urls = await page.evaluate('''() => {
            const urls = {};
            document.querySelectorAll("[data-fileurl]").forEach(el => {
                const href = el.getAttribute("data-fileurl") || "";
                const match = href.match(/\\/(R\\d+)_/);
                if (match) {
                    const key = match[1];
                    if (!urls[key]) urls[key] = href;
                }
            });
            return urls;
        }''')
        file_urls.update(dom_file_urls)

    except Exception as e:
        print(f"  페이지 로드 오류: {e}")
    finally:
        page.remove_listener("response", on_response)

    return api_calls, file_urls


def download_pdf(pdf_url, save_path, session_cookies=None):
    """PDF 다운로드"""
    headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 Chrome/120.0.0.0 Safari/537.36'}
    try:
        s = requests.Session()
        s.headers.update(headers)
        if session_cookies:
            for c in session_cookies:
                s.cookies.set(c['name'], c['value'], domain=c.get('domain', '.funetf.co.kr'))
        r = s.get(pdf_url, timeout=30)
        if r.status_code == 200 and len(r.content) > 1000:
            save_path.write_bytes(r.content)
            return True
        else:
            print(f"  PDF 다운로드 실패: {r.status_code}, url={pdf_url}")
            return False
    except Exception as e:
        print(f"  PDF 다운로드 오류: {e}")
        return False


def find_fee_table_in_pdf(pdf_path):
    """
    PDF에서 '보수 및 비용' 표를 찾아 페이지 번호와 표 위치 반환
    Returns: (page_num, table_info) or None
    """
    try:
        with pdfplumber.open(pdf_path) as pdf:
            target_keywords = ["집합투자기구에 부과되는 보수", "나. 집합투자기구", "총보수", "증권거래비용"]
            for page_num, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                if any(kw in text for kw in target_keywords[:2]):
                    # 이 페이지에 보수 표 있음
                    tables = page.extract_tables()
                    # 총보수나 증권거래비용이 포함된 표 찾기
                    for tbl in tables:
                        flat = " ".join(str(cell) for row in tbl for cell in row if cell)
                        if "총보수" in flat or "증권거래비용" in flat:
                            return page_num, tbl
                    # 표를 못 찾으면 텍스트에서 값 추출 시도
                    return page_num, None
        return None, None
    except Exception as e:
        print(f"  PDF 파싱 오류: {e}")
        return None, None


def extract_fee_values_from_pdf(pdf_path):
    """
    PDF에서 합성총보수(총보수)와 증권거래비용 수치 추출
    Returns: dict with keys 'total_fee', 'trading_cost', 'page_num'
    """
    result = {"total_fee": None, "trading_cost": None, "page_num": None, "all_classes": []}
    try:
        with pdfplumber.open(pdf_path) as pdf:
            fee_keywords = ["총보수", "증권거래비용", "집합투자기구에 부과되는 보수"]
            for page_num, page in enumerate(pdf.pages):
                text = page.extract_text() or ""

                # 보수 표 페이지 감지
                has_fee_section = any(kw in text for kw in ["집합투자기구에 부과되는 보수", "나. 집합투자기구에 부과"])
                has_fee_keywords = ("총보수" in text or "증권거래비용" in text)

                if not (has_fee_section or has_fee_keywords):
                    continue

                result["page_num"] = page_num

                # 텍스트에서 수치 추출
                lines = text.split("\n")
                for i, line in enumerate(lines):
                    # 총보수 추출 (증권거래비용 전 줄의 퍼센트 값)
                    if "합성총보수" in line or ("총보수" in line and "증권거래비용" not in line):
                        # 같은 줄이나 다음 줄에서 숫자 찾기
                        numbers = re.findall(r'\d+\.\d+', line + " " + (lines[i+1] if i+1 < len(lines) else ""))
                        if numbers:
                            result["total_fee"] = float(numbers[0])

                    if "증권거래비용" in line:
                        numbers = re.findall(r'\d+\.\d+', line + " " + (lines[i+1] if i+1 < len(lines) else ""))
                        if numbers:
                            result["trading_cost"] = float(numbers[0])

                # 표 데이터 추출
                tables = page.extract_tables()
                for tbl in tables:
                    if not tbl:
                        continue
                    flat = " ".join(str(c) for row in tbl for c in row if c)
                    if "총보수" in flat or "증권거래비용" in flat:
                        result["fee_table"] = tbl
                        # 표에서 클래스별 보수 추출
                        for row in tbl:
                            row_str = [str(c).strip() if c else "" for c in row]
                            if "증권거래비용" in row_str[0] if row_str else False:
                                numbers = [c for c in row_str[1:] if re.match(r'\d+\.\d+', c.strip())]
                                if numbers:
                                    result["trading_cost"] = float(numbers[0])
                        break

                if result["page_num"] is not None:
                    break

    except Exception as e:
        print(f"  PDF 값 추출 오류: {e}")

    return result


def screenshot_fee_table_from_pdf(pdf_path, page_num, output_img_path, zoom=2.0):
    """
    PyMuPDF로 PDF 특정 페이지를 고화질 이미지로 저장
    """
    try:
        doc = fitz.open(str(pdf_path))

        # 보수 표 페이지부터 시작
        # 경우에 따라 표가 2페이지에 걸칠 수 있음
        target_pages = [page_num]

        # 다음 페이지도 "보수" 키워드 있으면 포함
        if page_num + 1 < len(doc):
            next_page = doc[page_num + 1]
            next_text = next_page.get_text()
            if "증권거래비용" in next_text or "총보수" in next_text:
                target_pages.append(page_num + 1)

        images = []
        for pn in target_pages:
            page = doc[pn]
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat)
            img_bytes = pix.tobytes("png")
            img = Image.open(BytesIO(img_bytes))
            images.append(img)

        if len(images) == 1:
            images[0].save(str(output_img_path))
        else:
            # 세로로 합치기
            total_height = sum(img.height for img in images)
            max_width = max(img.width for img in images)
            combined = Image.new('RGB', (max_width, total_height), 'white')
            y = 0
            for img in images:
                combined.paste(img, (0, y))
                y += img.height
            combined.save(str(output_img_path))

        doc.close()
        return True
    except Exception as e:
        print(f"  PDF 스크린샷 오류: {e}")
        return False


def crop_fee_section_from_image(img_path, pdf_path, page_num):
    """
    이미지에서 보수 표 섹션만 크롭
    pdfplumber로 표 위치 파악 후 크롭
    """
    try:
        with pdfplumber.open(pdf_path) as pdf:
            page = pdf.pages[page_num]
            page_height = page.height
            page_width = page.width

            # 표 찾기
            tables = page.find_tables()
            fee_table = None
            for tbl in tables:
                # 표 텍스트 확인
                extracted = tbl.extract()
                if extracted:
                    flat = " ".join(str(c) for row in extracted for c in row if c)
                    if "총보수" in flat or "증권거래비용" in flat:
                        fee_table = tbl
                        break

            if fee_table is None:
                return img_path  # 크롭 없이 원본 반환

            # 섹션 헤더부터 포함하기 위해 y 좌표 위로 확장
            bbox = fee_table.bbox  # (x0, top, x1, bottom) in PDF units

            # 이미지 크기
            img = Image.open(str(img_path))
            img_w, img_h = img.size

            # PDF 좌표 → 이미지 좌표 변환
            scale_x = img_w / page_width
            scale_y = img_h / page_height

            # 여유 공간 추가
            margin_top = 40  # 제목 포함
            margin_bottom = 20

            x0 = max(0, int(bbox[0] * scale_x) - 10)
            y0 = max(0, int(bbox[1] * scale_y) - margin_top)
            x1 = min(img_w, int(bbox[2] * scale_x) + 10)
            y1 = min(img_h, int(bbox[3] * scale_y) + margin_bottom)

            cropped = img.crop((x0, y0, x1, y1))

            # 크롭된 이미지 저장
            crop_path = str(img_path).replace(".png", "_cropped.png")
            cropped.save(crop_path)
            return crop_path

    except Exception as e:
        print(f"  이미지 크롭 오류: {e}")
        return str(img_path)


def create_output_pptx(results, output_path):
    """결과 데이터로 PPT 파일 생성"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # 완전 빈 레이아웃

    for idx, res in enumerate(results):
        slide = prs.slides.add_slide(blank_layout)

        # ── 배경: 흰색 (기본값)

        # ── 좌측 상단: 펀드/ETF명 ──
        name = res.get("name", "")
        name_box = slide.shapes.add_textbox(
            Inches(0.2), Inches(0.1), Inches(12.9), Inches(0.6)
        )
        tf = name_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = name
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

        # ── 기본 정보 (좌측) ──
        info_lines = []
        if res.get("listing_date"):
            label = "상장일" if res.get("item_type") == "etf" else "설정일"
            info_lines.append(f"{label}: {res['listing_date']}")
        if res.get("amc"):
            info_lines.append(f"자산운용사: {res['amc']}")
        if res.get("risk_grade"):
            info_lines.append(f"투자위험도: {res['risk_grade']}")
        if res.get("total_fee") is not None:
            info_lines.append(f"합성총보수: {res['total_fee']}%")
        if res.get("trading_cost") is not None:
            info_lines.append(f"증권거래비용: {res['trading_cost']}%")

        if info_lines:
            info_box = slide.shapes.add_textbox(
                Inches(0.2), Inches(0.75), Inches(3.5), Inches(2.5)
            )
            tf = info_box.text_frame
            tf.word_wrap = True
            for i, line in enumerate(info_lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                run = p.add_run()
                run.text = line
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                p.space_after = Pt(4)

        # ── 보수 표 이미지 (우측) ──
        img_path = res.get("fee_image_path")
        if img_path and os.path.exists(img_path):
            try:
                img = Image.open(img_path)
                img_w, img_h = img.size
                aspect = img_h / img_w

                # 슬라이드 오른쪽에 배치
                pic_width = Inches(9.0)
                pic_height = pic_width * aspect

                # 높이가 너무 크면 조정
                max_height = Inches(6.5)
                if pic_height > max_height:
                    pic_height = max_height
                    pic_width = pic_height / aspect

                left = Inches(13.33) - pic_width - Inches(0.1)
                top = Inches(0.8)

                slide.shapes.add_picture(str(img_path), left, top, pic_width, pic_height)
            except Exception as e:
                print(f"  이미지 삽입 오류 ({name}): {e}")
        else:
            # 이미지 없을 때 텍스트로 표시
            msg_box = slide.shapes.add_textbox(Inches(4), Inches(2), Inches(8), Inches(1))
            tf = msg_box.text_frame
            p = tf.paragraphs[0]
            p.add_run().text = "PDF 이미지 없음"

        # ── 슬라이드 번호 ──
        num_box = slide.shapes.add_textbox(Inches(12.8), Inches(7.1), Inches(0.5), Inches(0.3))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = str(idx + 1)
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)

    prs.save(str(output_path))
    print(f"✅ PPT 저장: {output_path}")


def create_summary_pptx(results, output_path):
    """요약 데이터 표를 포함한 PPT 슬라이드 생성"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "2026년 5월 연령별 연금 포트폴리오 - 보수 현황"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # 표 추가
    rows = len(results) + 1
    cols = 6
    table = slide.shapes.add_table(
        rows, cols,
        Inches(0.2), Inches(0.7),
        Inches(12.9), Inches(6.5)
    ).table

    # 헤더
    headers = ["구분", "상장일/설정일", "자산운용사", "투자위험도", "합성총보수(%)", "증권거래비용(%)"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1F, 0x49, 0x7D)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 데이터 행
    for row_idx, res in enumerate(results):
        data = [
            res.get("name", "")[:40],
            res.get("listing_date", "-"),
            res.get("amc", "-"),
            res.get("risk_grade", "-"),
            str(res.get("total_fee", "-")),
            str(res.get("trading_cost", "-")),
        ]
        for col_idx, val in enumerate(data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(val) if val is not None else "-"
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFF)

    prs.save(str(output_path))
    print(f"✅ 요약 PPT 저장: {output_path}")


async def process_item(page, item, browser_cookies=None):
    """개별 항목 처리"""
    name = item["name"]
    item_type = item["type"]
    print(f"\n{'='*60}")
    print(f"처리 중: {name}")

    result = {
        "name": name,
        "item_type": item_type,
        "listing_date": None,
        "amc": None,
        "risk_grade": None,
        "total_fee": None,
        "trading_cost": None,
        "fee_image_path": None,
    }

    # ── 1. 검색 API로 기본 정보 수집 ──
    etf_list, fund_list = await search_fund_etf(page, item["search"])

    search_data = None
    fund_cd = None
    rep_fid = None

    if item_type == "etf":
        sot_cd = item["sotCd"]
        # sotCd로 매칭
        for etf in etf_list:
            if str(etf.get("sotCd", "")) == sot_cd:
                search_data = etf
                fund_cd = etf.get("fundCd")
                rep_fid = etf.get("repFId") or etf.get("fid")
                break
        if not search_data and etf_list:
            # sotCd 매칭 실패 시 첫 번째 결과 사용
            search_data = etf_list[0]
            fund_cd = search_data.get("fundCd")
            rep_fid = search_data.get("repFId") or search_data.get("fid")
            print(f"  ⚠️ sotCd 매칭 실패, 첫 번째 결과 사용: {search_data.get('itemNm')}")
    else:
        # 펀드 검색
        code = item.get("code", "")
        for fund in fund_list:
            if str(fund.get("shortCd", "")).upper() == code.upper():
                search_data = fund
                fund_cd = fund.get("fundCd")
                rep_fid = fund.get("repFId") or fund.get("fid")
                break
        if not search_data and fund_list:
            # 코드 매칭 실패 시 이름으로 매칭
            search_name_key = item["search"].replace(" ", "").lower()
            for fund in fund_list:
                fund_name = fund.get("fundFnm", "").replace(" ", "").lower()
                if search_name_key[:10] in fund_name:
                    search_data = fund
                    fund_cd = fund.get("fundCd")
                    rep_fid = fund.get("repFId") or fund.get("fid")
                    break
        if not search_data and fund_list:
            search_data = fund_list[0]
            fund_cd = search_data.get("fundCd")
            rep_fid = search_data.get("repFId") or search_data.get("fid")
            print(f"  ⚠️ 코드 매칭 실패, 첫 번째 결과 사용: {search_data.get('fundFnm')}")

    if search_data:
        # 기본 정보 추출
        if item_type == "etf":
            result["listing_date"] = fmt_date(search_data.get("lstnDt") or search_data.get("seoljYmd"))
        else:
            result["listing_date"] = fmt_date(search_data.get("seoljYmd") or search_data.get("lstnDt"))

        result["amc"] = search_data.get("unyongNm", "-")
        result["risk_grade"] = risk_grade_label(search_data.get("investGrade"))

        # 합성총보수: ETF는 ter, 펀드는 feeTot (임시)
        ter = search_data.get("ter")
        fee_tot = search_data.get("feeTot")
        if ter is not None:
            result["total_fee"] = round(float(ter), 4)
        elif fee_tot is not None:
            result["total_fee"] = round(float(fee_tot), 4)

        print(f"  ✅ 검색 성공: {search_data.get('fundFnm') or search_data.get('itemNm')}")
        print(f"     fundCd={fund_cd}, repFId={rep_fid}")
        print(f"     설정일={result['listing_date']}, 운용사={result['amc']}, 위험도={result['risk_grade']}")
        print(f"     합성총보수={result['total_fee']}")
    else:
        print(f"  ❌ 검색 결과 없음")
        return result

    # ── 2. 상세 페이지 로드 → etfdoc API 캡쳐 ──
    if item_type == "etf":
        detail_url = f"https://www.funetf.co.kr/product/etf/view/{item['sotCd']}"
    else:
        detail_url = f"https://www.funetf.co.kr/product/fund/view/{fund_cd}"

    print(f"  📄 상세 페이지 로드: {detail_url}")
    api_data, file_urls = await get_detail_page_data(page, detail_url)

    # 투자설명서 URL 결정 (우선순위: DOM > etfdoc API > repFId 추정)
    pdf_url = None
    pdf_filename = None

    # 1. DOM에서 R1 파일 URL (투자설명서) 우선
    if "R1" in file_urls:
        r1_path = file_urls["R1"]
        pdf_url = f"https://www.funetf.co.kr{r1_path}"
        pdf_filename = r1_path.split("/")[-1]
        print(f"  📋 투자설명서 URL (DOM): {pdf_url}")
    elif "R2" in file_urls:
        # 간이 투자설명서 (차선책)
        r2_path = file_urls["R2"]
        pdf_url = f"https://www.funetf.co.kr{r2_path}"
        pdf_filename = r2_path.split("/")[-1]
        print(f"  📋 간이투자설명서 URL (DOM): {pdf_url}")

    # 2. etfdoc API 응답에서 파일명 (삼성 계열)
    if not pdf_url and "etfdoc" in api_data:
        doc_data = api_data["etfdoc"]
        content = doc_data.get("content", []) if isinstance(doc_data, dict) else doc_data
        if isinstance(content, list) and content:
            explain_t = content[0].get("explainT")
            if explain_t:
                pdf_filename = explain_t
                pdf_url = f"https://www.funetf.co.kr/upload/invest/{pdf_filename}"
                print(f"  📋 투자설명서 URL (etfdoc): {pdf_url}")

    # 3. repFId 기반 추정 (KODEX 계열)
    if not pdf_url and rep_fid:
        pdf_filename = f"{rep_fid}-A.pdf"
        pdf_url = f"https://www.funetf.co.kr/upload/invest/{pdf_filename}"
        print(f"  ⚠️ 추정 PDF URL: {pdf_url}")

    if not pdf_url:
        print(f"  ❌ PDF URL 확인 불가")
        return result
    pdf_save_path = PDF_DIR / pdf_filename

    if not pdf_save_path.exists():
        print(f"  📥 PDF 다운로드: {pdf_url}")
        success = download_pdf(pdf_url, pdf_save_path, browser_cookies)
        if not success:
            # R2 (간이) 실패시 R3 시도
            if "R3" in file_urls:
                r3_path = file_urls["R3"]
                alt_url = f"https://www.funetf.co.kr{r3_path}"
                alt_filename = r3_path.split("/")[-1]
                alt_path = PDF_DIR / alt_filename
                print(f"  🔄 핵심정보설명서 시도: {alt_url}")
                if download_pdf(alt_url, alt_path, browser_cookies):
                    pdf_save_path = alt_path
                    pdf_filename = alt_filename
                else:
                    print(f"  ❌ PDF 다운로드 실패")
                    return result
            else:
                print(f"  ❌ PDF 다운로드 실패")
                return result
    else:
        print(f"  ♻️ PDF 캐시 사용: {pdf_filename}")

    # ── 4. PDF에서 보수 정보 추출 ──
    print(f"  🔍 PDF에서 보수 정보 추출...")
    fee_info = extract_fee_values_from_pdf(pdf_save_path)

    if fee_info.get("total_fee") is not None:
        result["total_fee"] = round(fee_info["total_fee"], 4)
    if fee_info.get("trading_cost") is not None:
        result["trading_cost"] = round(fee_info["trading_cost"], 4)

    print(f"  💰 합성총보수={result['total_fee']}, 증권거래비용={result['trading_cost']}")

    # ── 5. 보수 표 페이지 스크린샷 ──
    page_num = fee_info.get("page_num")
    if page_num is not None:
        # 이미지 파일명: PDF이름_페이지.png
        safe_name = re.sub(r'[^\w\-_]', '_', name)[:40]
        img_filename = f"{safe_name}.png"
        img_path = IMG_DIR / img_filename

        print(f"  📸 보수 표 스크린샷 (페이지 {page_num+1})...")
        if screenshot_fee_table_from_pdf(pdf_save_path, page_num, img_path):
            # 표 영역만 크롭
            crop_path = crop_fee_section_from_image(img_path, pdf_save_path, page_num)
            result["fee_image_path"] = crop_path
            print(f"  ✅ 이미지 저장: {crop_path}")
        else:
            print(f"  ❌ 스크린샷 실패")

    return result


async def main():
    all_results = []

    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        context = await browser.new_context(
            user_agent='Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 Chrome/120.0.0.0 Safari/537.36'
        )
        page = await context.new_page()

        # 로그인
        await login(page)

        # 브라우저 쿠키 가져오기 (PDF 다운로드용)
        cookies = await context.cookies()

        # 각 항목 처리
        for idx, item in enumerate(ITEMS):
            try:
                result = await process_item(page, item, browser_cookies=cookies)
                all_results.append(result)

                # 진행률 출력
                print(f"\n  진행률: {idx+1}/{len(ITEMS)}")

            except Exception as e:
                print(f"\n❌ 오류 ({item['name']}): {e}")
                import traceback
                traceback.print_exc()
                all_results.append({
                    "name": item["name"],
                    "item_type": item["type"],
                    "listing_date": None,
                    "amc": None,
                    "risk_grade": None,
                    "total_fee": None,
                    "trading_cost": None,
                    "fee_image_path": None,
                    "error": str(e),
                })

        await browser.close()

    # ── 결과 출력 ──
    print("\n" + "="*70)
    print("📊 수집 결과 요약")
    print("="*70)
    print(f"{'항목명':<45} {'상장/설정일':<12} {'운용사':<10} {'위험도':<8} {'합성총보수':<10} {'증권거래비용'}")
    print("-"*110)
    for res in all_results:
        print(f"{res['name'][:44]:<45} {str(res.get('listing_date') or '-'):<12} {str(res.get('amc') or '-'):<10} {str(res.get('risk_grade') or '-')[:7]:<8} {str(res.get('total_fee') or '-'):<10} {str(res.get('trading_cost') or '-')}")

    # ── PPT 생성 ──
    timestamp = datetime.now().strftime("%Y%m%d_%H%M")

    # 보수 표 캡쳐 PPT
    pptx_path = OUTPUT_DIR / f"보수표_캡쳐_{timestamp}.pptx"
    create_output_pptx(all_results, pptx_path)

    # 요약 데이터 PPT
    summary_path = OUTPUT_DIR / f"보수_요약_{timestamp}.pptx"
    create_summary_pptx(all_results, summary_path)

    # JSON 결과 저장
    json_path = OUTPUT_DIR / f"results_{timestamp}.json"
    with open(json_path, 'w', encoding='utf-8') as f:
        json.dump(all_results, f, ensure_ascii=False, indent=2, default=str)
    print(f"✅ JSON 저장: {json_path}")

    print("\n🎉 완료!")


if __name__ == "__main__":
    asyncio.run(main())
